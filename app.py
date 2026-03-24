# ── Ensure db/ and cache/ directories exist at their configured locations ──
# Must run BEFORE importing config/database/file_index/file_monitor so that
# those modules find their directories already in place.
# setup_storage.py and config.py intentionally do NOT call this.
from paths import ensure_dirs

ensure_dirs()

# Bulk ZIP progress tracking
bulk_zip_progress = {}

from flask import g

# Global flag for cancelling bulk ZIP
bulk_zip_cancelled = {}

# Move this endpoint below app initialization
from flask import (
    Flask,
    render_template,
    request,
    redirect,
    url_for,
    send_from_directory,
    send_file,
    flash,
    session,
    jsonify,
    Response,
    make_response,
    render_template_string,
)
from flask_cors import CORS
from werkzeug.utils import secure_filename
from werkzeug.exceptions import ClientDisconnected
import os
import sys
import shutil
import json
import threading
import time
import logging
import uuid

# Ensure ffmpeg and thread prints appear immediately in terminal
sys.stdout.reconfigure(line_buffering=True)
import zipfile
import io
import re
import subprocess
import hashlib
import zipstream
from datetime import datetime
from config import (
    PORT,
    ROOT_DIR,
    CHUNK_SIZE,
    ENABLE_CHUNKED_UPLOADS,
    HLS_MIN_SIZE,
    HLS_FORCE_FORMATS,
)
from database import get_session_secret

SESSION_SECRET = get_session_secret()
from auth import (
    check_login,
    login_user,
    logout_user,
    current_user,
    is_logged_in,
    get_role,
)
import storage


# ------------------------------------------------------------------
# Brute-force protection — tracks failed login attempts per IP
# ------------------------------------------------------------------
class RateLimiter:
    """
    Tracks failed login attempts per IP address.
    After MAX_ATTEMPTS failures within WINDOW seconds, the IP is locked
    out for LOCKOUT seconds.
    All state is in-memory — resets on server restart (intentional).
    """

    MAX_ATTEMPTS = 5  # failures before lockout
    WINDOW = 60  # seconds — rolling window for counting failures
    LOCKOUT = 300  # seconds — how long the IP is blocked (5 minutes)

    def __init__(self):
        self._attempts: dict = {}  # ip -> list of failure timestamps
        self._locked: dict = {}  # ip -> lockout-expiry timestamp
        self._lock = threading.Lock()

    def _get_ip(self) -> str:
        return (
            request.headers.get("X-Forwarded-For", request.remote_addr or "")
            .split(",")[0]
            .strip()
        )

    def is_blocked(self) -> bool:
        ip = self._get_ip()
        with self._lock:
            expiry = self._locked.get(ip)
            if expiry:
                if time.time() < expiry:
                    return True
                else:
                    del self._locked[ip]
                    self._attempts.pop(ip, None)
        return False

    def record_failure(self):
        ip = self._get_ip()
        now = time.time()
        with self._lock:
            attempts = [t for t in self._attempts.get(ip, []) if now - t < self.WINDOW]
            attempts.append(now)
            self._attempts[ip] = attempts
            if len(attempts) >= self.MAX_ATTEMPTS:
                self._locked[ip] = now + self.LOCKOUT
                print(f"IP locked out after {self.MAX_ATTEMPTS} failed attempts: {ip}")

    def record_success(self):
        ip = self._get_ip()
        with self._lock:
            self._attempts.pop(ip, None)
            self._locked.pop(ip, None)

    def remaining_lockout(self) -> int:
        ip = self._get_ip()
        with self._lock:
            expiry = self._locked.get(ip, 0)
            return max(0, int(expiry - time.time()))

    def attempts_remaining(self) -> int:
        ip = self._get_ip()
        now = time.time()
        with self._lock:
            recent = [t for t in self._attempts.get(ip, []) if now - t < self.WINDOW]
            return max(0, self.MAX_ATTEMPTS - len(recent))


rate_limiter = RateLimiter()

# File monitoring and real-time updates
from file_monitor import get_file_monitor, init_file_monitor
from realtime_stats import storage_stats_sse, trigger_storage_update, get_event_manager

# Assembly Queue System
import queue
from dataclasses import dataclass
from typing import Optional


@dataclass
class AssemblyJob:
    file_id: str
    filename: str
    dest_path: str
    total_chunks: int
    created_at: float
    status: str = "pending"  # pending, processing, completed, error
    error_message: Optional[str] = None
    session_id: Optional[str] = None


class AssemblyQueue:
    def __init__(self):
        self.job_queue = queue.Queue()
        self.active_jobs = {}  # file_id -> AssemblyJob
        self.completed_jobs = {}  # file_id -> AssemblyJob (keep for 1 hour)
        self.lock = threading.Lock()

    def add_job(self, file_id, filename, dest_path, total_chunks, session_id=None):
        """Add a new assembly job to the queue"""
        job = AssemblyJob(
            file_id=file_id,
            filename=filename,
            dest_path=dest_path,
            total_chunks=total_chunks,
            created_at=time.time(),
            session_id=session_id,
        )

        with self.lock:
            self.active_jobs[file_id] = job

        self.job_queue.put(job)
        print(f"🔄 Added assembly job for {filename} (ID: {file_id})")
        return job

    def get_job_status(self, file_id):
        """Get the current status of an assembly job"""
        with self.lock:
            if file_id in self.active_jobs:
                return self.active_jobs[file_id]
            elif file_id in self.completed_jobs:
                return self.completed_jobs[file_id]
            return None

    def get_all_active_jobs(self):
        """Get all currently active assembly jobs"""
        with self.lock:
            return list(self.active_jobs.values())

    def get_jobs_for_session(self, session_id):
        """Get all jobs (active + recent completed) for a session"""
        with self.lock:
            jobs = []
            # Active jobs
            for job in self.active_jobs.values():
                if job.session_id == session_id:
                    jobs.append(job)
            # Recent completed jobs (last hour)
            current_time = time.time()
            for job in self.completed_jobs.values():
                if (
                    job.session_id == session_id
                    and current_time - job.created_at < 3600
                ):  # 1 hour
                    jobs.append(job)
            return jobs

    def complete_job(self, file_id, success=True, error_message=None):
        """Mark a job as completed"""
        with self.lock:
            if file_id in self.active_jobs:
                job = self.active_jobs.pop(file_id)
                job.status = "completed" if success else "error"
                if error_message:
                    job.error_message = error_message
                self.completed_jobs[file_id] = job
                print(
                    f"✅ Assembly job completed for {job.filename} (Success: {success})"
                )

                # Untrack the upload when assembly is successfully completed
                if success and job.session_id:
                    try:
                        chunk_tracker.untrack_upload(job.session_id, file_id)
                        print(
                            f"🧹 Untracked completed upload: {file_id} for session {job.session_id}"
                        )
                    except Exception as e:
                        print(f"⚠️ Failed to untrack upload {file_id}: {e}")

                return job
            return None

    def cleanup_old_jobs(self):
        """Remove completed jobs older than 1 hour"""
        with self.lock:
            current_time = time.time()
            expired_jobs = [
                file_id
                for file_id, job in self.completed_jobs.items()
                if current_time - job.created_at > 3600
            ]
            for file_id in expired_jobs:
                del self.completed_jobs[file_id]
            if expired_jobs:
                print(f"🧹 Cleaned up {len(expired_jobs)} old assembly jobs")


# Global assembly queue
assembly_queue = AssemblyQueue()

app = Flask(__name__)
CORS(app)
app.secret_key = SESSION_SECRET

# Configure session handling
app.config.update(
    SESSION_COOKIE_SAMESITE="Lax",
    PERMANENT_SESSION_LIFETIME=86400,  # 24 hours
    SESSION_REFRESH_EACH_REQUEST=True,
    SESSION_COOKIE_NAME="cloudinator_session",
)


def _trigger_reconcile(settle=False):
    """Kick off a background reconcile so file/dir counts correct themselves
    immediately after mutations (delete, move, rename, copy) instead of waiting 15 min.

    IMPORTANT: We capture the snapshot BEFORE reconcile starts and always force-push
    SSE afterwards. Without this, a race between watchdog and the reconcile walk
    causes both to see no change and neither fires the SSE update:

      1. Copy finishes -> watchdog on_created events update _file_count/_total_size
         in memory but have NOT yet called _notify_and_save / updated last_snapshot.
      2. Reconcile walk finishes -> old_snapshot = build_snapshot() already has the
         NEW correct counts (watchdog updated the counters) -> old == new ->
         _notify_changes NOT called -> last_snapshot set to new correct values.
      3. Watchdog debounce fires -> old = last_snapshot (now the new correct values
         set by reconcile) -> new = build_snapshot() (same) -> no push again.
      Result: zero SSE pushes, client stuck until manual refresh.

    Fix: _reconcile() now always force-pushes a reconcile_complete=True SSE at the
    end regardless of drift, so this explicit trigger_storage_update call below is
    a belt-and-suspenders safety net only.

    NOTE: settle=True used to arm set_pending_reconcile() AFTER the walk — that
    caused an unnecessary 4th full walk per copy.  It is no longer needed because
    _reconcile() already handles the epoch guard that prevents double-counting from
    the watchdog backlog that drains after copytree finishes.
    """
    import threading
    from file_monitor import get_file_monitor
    from realtime_stats import trigger_storage_update

    def _run():
        try:
            monitor = get_file_monitor()
            # _reconcile() will push incremental walk-progress SSE + a final
            # reconcile_complete SSE internally.  The explicit push below is
            # kept only as a safety net for the case where the walk saw no drift
            # (old == new) but the race condition described above means last_snapshot
            # was already updated before reconcile ran.
            monitor._reconcile()
            new_snap = monitor.get_current_snapshot()
            print(
                f"\U0001f4e1 Reconcile complete (settle={settle}): "
                f"{getattr(new_snap, 'file_count', '?')} files"
            )
        except Exception as e:
            print(f"\u26a0\ufe0f Background reconcile error: {e}")

    threading.Thread(target=_run, daemon=True).start()


storage.ensure_root()

# Initialize file system monitoring
file_monitor = init_file_monitor()
file_monitor.add_change_callback(trigger_storage_update)
print(f"📡 File system monitoring started for: {ROOT_DIR}")


@app.before_request
def validate_session():
    # Skip validation for login-related routes
    if request.endpoint in ["login", "static"]:
        return

    # Check if user is logged in
    if not session.get("logged_in"):
        session.clear()
        return redirect(url_for("login"))

    # Verify the account still exists in users.json.
    # Without this, a deleted account's still-valid cookie causes an infinite
    # loop: /admin/upload_status returns 403 (role=None) → JS redirects to
    # /login → /login sees logged_in=True → redirects back to / → repeat.
    username = session.get("username")
    if not username or get_role(username) is None:
        session.clear()
        return redirect(url_for("login"))

    # Session lifetime controlled by PERMANENT_SESSION_LIFETIME (86400s = 24h)
    # and refreshed on every request via SESSION_REFRESH_EACH_REQUEST=True.


@app.route("/check_session")
def check_session():
    if not is_logged_in():
        return jsonify({"error": "Session expired"}), 401
    return jsonify({"status": "ok"}), 200


@app.route("/cancel_bulk_zip", methods=["POST"])
def cancel_bulk_zip():
    session_id = session.get("session_id") or request.cookies.get("session")
    if not session_id:
        return jsonify({"error": "No session ID"}), 400
    bulk_zip_cancelled[session_id] = True
    print(f"❌ Bulk ZIP cancelled for session {session_id}")
    return jsonify({"status": "cancelled"})


# Add Jinja2 filter for timestamp formatting
@app.template_filter("timestamp_to_date")
def timestamp_to_date_filter(timestamp):
    """Convert Unix timestamp to time on first line, date on second"""
    try:
        dt = datetime.fromtimestamp(timestamp)
        return dt.strftime("%m/%d/%Y") + "||" + dt.strftime("%I:%M %p")
    except (ValueError, OSError):
        return "--"


# Chunk Tracker Class for better session management
class ChunkTracker:
    def __init__(self):
        self.active_uploads = {}  # session_id -> set of file_ids
        self.lock = threading.Lock()
        self.upload_timestamps = {}  # file_id -> timestamp

    def track_upload(self, session_id, file_id):
        with self.lock:
            if session_id not in self.active_uploads:
                self.active_uploads[session_id] = set()
            self.active_uploads[session_id].add(file_id)
            self.upload_timestamps[file_id] = time.time()
            print(f"📊 Tracking upload: {file_id} for session {session_id}")

    def untrack_upload(self, session_id, file_id):
        with self.lock:
            if session_id in self.active_uploads:
                self.active_uploads[session_id].discard(file_id)
                if not self.active_uploads[session_id]:
                    del self.active_uploads[session_id]
            self.upload_timestamps.pop(file_id, None)
            print(f"📊 Untracked upload: {file_id} for session {session_id}")

    def cleanup_session_chunks(self, session_id):
        """Clean up all chunks for a session"""
        with self.lock:
            if session_id in self.active_uploads:
                file_ids = self.active_uploads[session_id].copy()
                for file_id in file_ids:
                    try:
                        storage.cleanup_chunks(file_id)
                        print(
                            f"🧹 Cleaned up abandoned chunks for session {session_id}: {file_id}"
                        )
                    except Exception as e:
                        print(f"❌ Error cleaning up chunks for {file_id}: {e}")
                    self.upload_timestamps.pop(file_id, None)
                del self.active_uploads[session_id]
                print(f"🧹 Cleaned up all chunks for session: {session_id}")

    def cleanup_orphaned_chunks(self):
        """Find and cleanup chunks that don't belong to any active session"""
        chunks_dir = os.path.join(ROOT_DIR, ".chunks")
        if not os.path.exists(chunks_dir):
            return

        try:
            with self.lock:
                all_tracked_files = set()
                for file_ids in self.active_uploads.values():
                    all_tracked_files.update(file_ids)

                current_time = time.time()
                cleaned_count = 0

                # Find chunks that aren't tracked by any session
                for file_id in os.listdir(chunks_dir):
                    chunk_dir = os.path.join(chunks_dir, file_id)
                    if not os.path.isdir(chunk_dir):
                        continue

                    # CRITICAL: Never cleanup chunks for files currently being assembled
                    if assembly_queue.get_job_status(file_id):
                        print(
                            f"🔐 Skipping cleanup for {file_id} - currently being assembled"
                        )
                        continue

                    # Also check for assembly protection marker
                    assembly_marker = os.path.join(chunk_dir, ".assembling")
                    if os.path.exists(assembly_marker):
                        print(
                            f"🔐 Skipping cleanup for {file_id} - assembly marker present"
                        )
                        continue

                    should_cleanup = False
                    cleanup_reason = ""

                    if file_id not in all_tracked_files:
                        # Check age before cleanup
                        timestamp_file = os.path.join(chunk_dir, ".timestamp")

                        if os.path.exists(timestamp_file):
                            try:
                                with open(timestamp_file, "r") as f:
                                    timestamp = float(f.read().strip())
                                # Cleanup untracked chunks older than 45 minutes.
                                # Must be > cleanup_interrupted_uploads timeout (30 min)
                                # so we never delete chunks that are simply backgrounded.
                                if current_time - timestamp > 2700:
                                    should_cleanup = True
                                    cleanup_reason = (
                                        f"untracked >45min old (abandoned upload)"
                                    )
                            except (ValueError, OSError):
                                should_cleanup = True
                                cleanup_reason = "corrupted timestamp file"
                        else:
                            # No timestamp, cleanup if dir is older than 45 minutes
                            try:
                                dir_mtime = os.path.getmtime(chunk_dir)
                                if current_time - dir_mtime > 2700:
                                    should_cleanup = True
                                    cleanup_reason = "no timestamp >45min old"
                            except OSError:
                                should_cleanup = True
                                cleanup_reason = "cannot read metadata"
                    else:
                        # Even tracked files - cleanup if very old (stale uploads)
                        file_timestamp = self.upload_timestamps.get(
                            file_id, current_time
                        )
                        if current_time - file_timestamp > 3600:  # 1 hour
                            should_cleanup = True
                            cleanup_reason = "tracked but stale >1hr"
                            print(
                                f"🧹 Cleaning up stale tracked chunks (>1hr): {file_id}"
                            )

                    if should_cleanup:
                        try:
                            storage.cleanup_chunks(file_id)
                            cleaned_count += 1
                            print(
                                f"🧹 Cleaned up orphaned chunks: {file_id} ({cleanup_reason})"
                            )
                        except Exception as e:
                            print(
                                f"❌ Failed to cleanup orphaned chunks {file_id}: {e}"
                            )

                        # Remove from tracking if it was tracked
                        if file_id in all_tracked_files:
                            for session_id, file_set in self.active_uploads.items():
                                file_set.discard(file_id)
                            self.upload_timestamps.pop(file_id, None)

                if cleaned_count > 0:
                    print(
                        f"🧹 Orphaned chunk cleanup completed: {cleaned_count} directories removed"
                    )

        except Exception as e:
            print(f"❌ Error in orphaned chunk cleanup: {e}")

    def cleanup_interrupted_uploads(self):
        """Detect and cleanup uploads that were interrupted (no activity for 30+ minutes).

        NOTE: The timeout is intentionally long (30 min) so that background-tab throttling
        does NOT trigger premature cleanup.  Browsers freeze JS timers/fetch in hidden tabs,
        so a 2-minute window incorrectly treated active uploads as abandoned whenever the
        user switched away for more than 2 minutes.  30 minutes gives plenty of headroom
        for large folder uploads that are paused while the user works in another tab.
        """
        current_time = time.time()
        interrupted_uploads = []

        try:
            with self.lock:
                for session_id, file_ids in list(self.active_uploads.items()):
                    for file_id in list(file_ids):
                        timestamp = self.upload_timestamps.get(file_id)
                        if (
                            timestamp and (current_time - timestamp) > 1800
                        ):  # 30 minutes of inactivity
                            interrupted_uploads.append((session_id, file_id))
                            print(
                                f"🧹 Detected interrupted upload: {file_id} (inactive for {int(current_time - timestamp)}s)"
                            )

                # Clean up interrupted uploads
                for session_id, file_id in interrupted_uploads:
                    try:
                        self.untrack_upload(session_id, file_id)
                        storage.cleanup_chunks(file_id)
                        print(f"🧹 Cleaned up interrupted upload: {file_id}")
                    except Exception as e:
                        print(f"❌ Failed to cleanup interrupted upload {file_id}: {e}")

                if interrupted_uploads:
                    print(
                        f"🧹 Interrupted upload cleanup completed: {len(interrupted_uploads)} uploads cleaned"
                    )

        except Exception as e:
            print(f"❌ Error in interrupted upload cleanup: {e}")

    def get_stats(self):
        """Get statistics about active uploads"""
        with self.lock:
            total_sessions = len(self.active_uploads)
            total_uploads = sum(
                len(file_set) for file_set in self.active_uploads.values()
            )
            return {
                "active_sessions": total_sessions,
                "active_uploads": total_uploads,
                "tracked_files": list(self.upload_timestamps.keys()),
            }


# Global chunk tracker instance
chunk_tracker = ChunkTracker()


def login_required(f):
    from functools import wraps

    @wraps(f)
    def decorated(*args, **kwargs):
        if not is_logged_in():
            # API and XHR requests: return 401 JSON so the client can handle it
            # without breaking an in-progress upload with a page redirect.
            if (
                request.path.startswith("/api/")
                or request.headers.get("X-Requested-With") == "XMLHttpRequest"
                or request.accept_mimetypes.best == "application/json"
            ):
                return jsonify({"error": "Session expired", "redirect": "/login"}), 401
            return redirect(url_for("login"))
        return f(*args, **kwargs)

    return decorated


def get_protected_files():
    """Get set of file IDs that should be protected from cleanup (currently being assembled)"""
    protected_files = set()
    try:
        for job in assembly_queue.get_all_active_jobs():
            protected_files.add(job.file_id)
    except Exception as e:
        print(f"⚠️ Warning: Could not get active assembly jobs: {e}")

    return protected_files


def cleanup_stale_chunks_on_request():
    """Clean up chunks that are older than 1 hour - called on each request"""
    try:
        # Get all active assembly jobs to avoid cleaning their chunks
        active_assembly_jobs = get_protected_files()

        if active_assembly_jobs:
            print(
                f"🔐 Protecting {len(active_assembly_jobs)} files from cleanup (currently being assembled)"
            )

        # Pass the protected file IDs to the cleanup function
        storage.cleanup_old_chunks(
            max_age_hours=1, protected_files=active_assembly_jobs
        )
    except Exception as e:
        print(f"❌ Error in stale chunk cleanup: {e}")


@app.before_request
def before_request():
    """Run cleanup before certain requests and handle interrupted uploads"""
    # Ensure session ID exists for logged-in users
    if is_logged_in() and "session_id" not in session:
        session["session_id"] = str(uuid.uuid4())

    # Enhanced cleanup on page load/refresh - check for assembly jobs first
    if request.endpoint in ["index", "upload"]:
        # Check for and cleanup any stale uploads from this session
        session_id = session.get("session_id")
        if session_id and request.endpoint == "index":
            # This is a page load/refresh - check for abandoned uploads
            # IMPORTANT: Check for assembly jobs FIRST before cleaning up chunks
            try:
                current_uploads = chunk_tracker.active_uploads.get(session_id, set())
                if current_uploads:
                    print(
                        f"🧹 Detected {len(current_uploads)} potentially abandoned uploads on page refresh"
                    )
                    # Check if any of these uploads are actually in assembly queue
                    assembly_protected = set()
                    for file_id in current_uploads.copy():
                        # Check if this file is in assembly queue
                        if assembly_queue.get_job_status(file_id):
                            print(f"🔐 Upload {file_id} is protected by assembly queue")
                            assembly_protected.add(file_id)
                            continue

                        # Check for assembly protection marker
                        chunk_dir = os.path.join(ROOT_DIR, ".chunks", file_id)
                        assembly_marker = os.path.join(chunk_dir, ".assembling")
                        if os.path.exists(assembly_marker):
                            print(
                                f"🔐 Upload {file_id} is protected by assembly marker"
                            )
                            assembly_protected.add(file_id)
                            continue

                        # Give a grace period for genuine page refreshes during upload
                        timestamp = chunk_tracker.upload_timestamps.get(file_id)
                        if (
                            timestamp and (time.time() - timestamp) > 30
                        ):  # 30 seconds grace period
                            print(f"🧹 Cleaning up abandoned upload: {file_id}")
                            chunk_tracker.untrack_upload(session_id, file_id)
                            storage.cleanup_chunks(file_id)

                    # Keep assembly-protected uploads in tracker
                    if assembly_protected:
                        print(
                            f"🔐 Keeping {len(assembly_protected)} assembly-protected uploads in tracker"
                        )
            except Exception as e:
                print(f"❌ Error in abandoned upload cleanup: {e}")

        # Run periodic cleanup in background thread to not slow down requests
        cleanup_thread = threading.Thread(
            target=cleanup_stale_chunks_on_request, daemon=True
        )
        cleanup_thread.start()


@app.after_request
def after_request(response):
    """Add security headers to all responses"""
    # Add cache control headers to authenticated pages
    if request.endpoint and request.endpoint not in ["login", "static"]:
        # Check if this is an authenticated route
        if is_logged_in() or request.endpoint in [
            "index",
            "download",
            "upload",
            "admin",
        ]:
            response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
            response.headers["Pragma"] = "no-cache"
            response.headers["Expires"] = "0"

    return response


@app.route("/login", methods=["GET", "POST"])
def login():
    # If user is already logged in, redirect to index
    if session.get("logged_in"):
        return redirect(url_for("index"))

    if request.method == "POST":
        # Check brute-force lockout before touching DB
        if rate_limiter.is_blocked():
            remaining = rate_limiter.remaining_lockout()
            flash(f"Too many failed attempts. Try again in {remaining} seconds.")
            return render_template("login.html"), 429

        username = request.form.get("username", "").strip()
        password = request.form.get("password", "")

        if check_login(username, password):
            rate_limiter.record_success()
            # Set up session data
            session.clear()
            session.permanent = True
            login_user(username)
            session["role"] = get_role(username)
            session["session_id"] = str(uuid.uuid4())
            session["logged_in"] = True
            session["login_time"] = int(time.time())
            session.modified = True

            return redirect(url_for("index"))
        else:
            rate_limiter.record_failure()
            left = rate_limiter.attempts_remaining()
            if left > 0:
                flash(f"Invalid username or password. {left} attempt(s) remaining.")
            else:
                flash(
                    f"Too many failed attempts. Try again in {RateLimiter.LOCKOUT} seconds."
                )
            return render_template("login.html"), 401

    # Render login page with no-cache headers
    response = make_response(render_template("login.html"))
    response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate, max-age=0"
    response.headers["Pragma"] = "no-cache"
    response.headers["Expires"] = "0"
    return response


@app.route("/logout")
def logout():
    try:
        # Clean up any upload chunks
        session_id = session.get("session_id")
        if session_id:
            chunk_tracker.cleanup_session_chunks(session_id)

        # Clear the session completely
        session.clear()

        # Create response with session-clearing headers
        response = make_response(redirect(url_for("login", logged_out="1")))
        response.delete_cookie("cloudinator_session")
        response.delete_cookie("session_check")

        # Add cache-control headers to prevent caching
        response.headers.update(
            {
                "Cache-Control": "no-cache, no-store, must-revalidate",
                "Pragma": "no-cache",
                "Expires": "0",
            }
        )

        return response
    except Exception as e:
        logging.error(f"Logout error: {e}", exc_info=True)
        session.clear()  # Still try to clear session even if other operations fail
        return redirect(url_for("login", logged_out="1"))


@app.route("/", defaults={"path": ""})
@app.route("/<path:path>")
@login_required
def index(path):
    # Comprehensive path validation: safety and existence
    if path and not storage.is_valid_path(path):
        if not storage.is_safe_path(path):
            flash("Invalid path: contains unsafe characters or directory traversal")
        else:
            flash(f'Path "{path}" does not exist or is not a directory')
        return redirect(url_for("index"))

    try:
        # Get current directory info
        current_path = os.path.join(ROOT_DIR, path) if path else ROOT_DIR
        items = storage.list_dir(path)

        response = make_response(
            render_template(
                "index.html",
                items=items,
                path=path,
                role=session.get("role", "readonly"),
                CHUNK_SIZE=CHUNK_SIZE,
            )
        )

        # Add strict cache control headers to prevent caching of authenticated content
        response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
        response.headers["Pragma"] = "no-cache"
        response.headers["Expires"] = "0"

        return response

    except Exception as e:
        logging.error(f"Error loading directory {path}: {e}", exc_info=True)
        flash("Error loading directory")
        return redirect(url_for("index"))


@app.route("/download/<path:path>")
@login_required
def download(path):
    # Security check: ensure path is safe
    if not storage.is_safe_path(path):
        flash("Invalid file path")
        return redirect(url_for("index"))

    full_path = os.path.join(ROOT_DIR, path)
    if not os.path.exists(full_path) or os.path.isdir(full_path):
        flash("File not found")
        return redirect(url_for("index"))

    directory = os.path.dirname(full_path)
    filename = os.path.basename(full_path)
    return send_from_directory(directory, filename, as_attachment=True)


@app.route("/view/<path:path>")
@login_required
def view_file(path):
    """Serve a file inline (for in-browser preview — images, video, audio, PDF, text)."""
    if not storage.is_safe_path(path):
        return "Invalid file path", 400
    full_path = os.path.join(ROOT_DIR, path)
    if not os.path.exists(full_path) or os.path.isdir(full_path):
        return "File not found", 404
    directory = os.path.dirname(full_path)
    filename = os.path.basename(full_path)
    return send_from_directory(directory, filename, as_attachment=False)


@app.route("/office_preview/<path:path>")
@login_required
def office_preview(path):
    """Convert Office documents (docx, xlsx, pptx) to structured data for browser preview."""
    import html as html_lib

    if not storage.is_safe_path(path):
        return jsonify({"error": "Invalid file path"}), 400
    full_path = os.path.join(ROOT_DIR, path)
    if not os.path.exists(full_path) or os.path.isdir(full_path):
        return jsonify({"error": "File not found"}), 404

    ext = path.rsplit(".", 1)[-1].lower() if "." in path else ""

    try:
        # ── DOCX ──────────────────────────────────────────────────────────────
        if ext in ("docx", "doc"):
            import mammoth

            with open(full_path, "rb") as f:
                result = mammoth.convert_to_html(f)
            return jsonify({"type": "docx", "html": result.value})

        # ── XLSX / XLS ────────────────────────────────────────────────────────
        elif ext in ("xlsx", "xls"):
            import openpyxl

            MAX_ROWS = 500
            MAX_COLS = 50
            wb = openpyxl.load_workbook(full_path, read_only=True, data_only=True)
            sheets = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for r_idx, row in enumerate(ws.iter_rows(values_only=True)):
                    if r_idx >= MAX_ROWS:
                        break
                    cells = [
                        html_lib.escape(str(c)) if c is not None else ""
                        for c in list(row)[:MAX_COLS]
                    ]
                    rows.append(cells)
                sheets.append(
                    {
                        "name": html_lib.escape(sheet_name),
                        "rows": rows,
                        "truncated": ws.max_row is not None and ws.max_row > MAX_ROWS,
                    }
                )
            wb.close()
            return jsonify({"type": "xlsx", "sheets": sheets})

        # ── CSV ───────────────────────────────────────────────────────────────────────
        elif ext == "csv":
            import csv as csv_mod

            MAX_ROWS = 500
            rows = []
            with open(full_path, newline="", encoding="utf-8", errors="replace") as f:
                reader = csv_mod.reader(f)
                for i, row in enumerate(reader):
                    if i >= MAX_ROWS:
                        break
                    rows.append([html_lib.escape(str(c)) for c in row])
            truncated = False
            try:
                # cheap line-count check without re-reading the whole file
                with open(full_path, encoding="utf-8", errors="replace") as f:
                    total_lines = sum(1 for _ in f)
                truncated = total_lines > MAX_ROWS
            except Exception:
                pass
            return jsonify(
                {
                    "type": "xlsx",  # reuse the same xlsx renderer on the frontend
                    "sheets": [{"name": "CSV", "rows": rows, "truncated": truncated}],
                }
            )

        # ── PPTX / PPT ────────────────────────────────────────────────────────
        elif ext in ("pptx", "ppt"):
            from pptx import Presentation

            try:
                from pptx.enum.shapes import PP_PLACEHOLDER

                _HAS_PLACEHOLDER_ENUM = True
            except ImportError:
                _HAS_PLACEHOLDER_ENUM = False

            prs = Presentation(full_path)
            slides = []
            for i, slide in enumerate(prs.slides):
                shapes_data = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    paragraphs = []
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        paragraphs.append(
                            {"text": html_lib.escape(text), "level": para.level}
                        )
                    if not paragraphs:
                        continue
                    is_title = False
                    if _HAS_PLACEHOLDER_ENUM:
                        try:
                            if (
                                shape.is_placeholder
                                and shape.placeholder_format.type
                                in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
                            ):
                                is_title = True
                        except Exception:
                            pass
                    shapes_data.append({"is_title": is_title, "paragraphs": paragraphs})
                slides.append({"index": i + 1, "shapes": shapes_data})
            return jsonify({"type": "pptx", "slides": slides, "total": len(slides)})

        else:
            return jsonify({"error": f"Unsupported format: .{ext}"}), 400

    except ImportError as e:
        missing = str(e).split("'")[1] if "'" in str(e) else str(e)
        return (
            jsonify(
                {
                    "error": f"Missing package '{missing}'. Install with: pip install mammoth openpyxl python-pptx"
                }
            ),
            500,
        )
    except Exception as e:
        logging.exception(f"Office preview failed for {path}")
        return jsonify({"error": f"Preview failed: {str(e)}"}), 500


@app.route("/archive_preview/<path:path>")
@login_required
def archive_preview(path):
    """List contents of an archive file for in-browser preview.

    Supports: .zip, .rar, .7z, .tar, .tar.gz, .tar.bz2, .tar.xz
    Optional query param: ?password=<password>

    Returns JSON:
      { type, entries: [{name, is_dir, size, compressed_size, modified}],
        total_entries, total_size, truncated, encrypted }

    Error responses (always JSON):
      {"error": "password_required"}   -> HTTP 401
      {"error": "wrong_password"}      -> HTTP 401
      {"error": "<message>"}           -> HTTP 400/500
    """
    import html as html_lib
    import datetime

    if not storage.is_safe_path(path):
        return jsonify({"error": "Invalid file path"}), 400
    full_path = os.path.join(ROOT_DIR, path)
    if not os.path.exists(full_path) or os.path.isdir(full_path):
        return jsonify({"error": "File not found"}), 404

    # Resolve compound extensions first (must check before splitting on last '.')
    lower = path.lower()
    if lower.endswith(".tar.gz") or lower.endswith(".tgz"):
        ext = "tar.gz"
    elif lower.endswith(".tar.bz2") or lower.endswith(".tbz2"):
        ext = "tar.bz2"
    elif lower.endswith(".tar.xz") or lower.endswith(".txz"):
        ext = "tar.xz"
    else:
        ext = lower.rsplit(".", 1)[-1] if "." in lower else ""

    password = request.args.get("password") or None
    MAX_ENTRIES = 10_000

    def _fmt_zip_date(dt_tuple):
        try:
            y, mo, d, h, mi = dt_tuple[:5]
            return f"{y:04d}-{mo:02d}-{d:02d} {h:02d}:{mi:02d}"
        except Exception:
            return ""

    try:
        # ── ZIP ───────────────────────────────────────────────────────────────
        if ext == "zip":
            # Use pyzipper when available — it handles WinZip AES encryption
            # (compression method 99). Python's stdlib zipfile cannot decrypt
            # AES zips: both correct and wrong passwords raise the same
            # NotImplementedError("That compression method is not supported"),
            # making password verification completely impossible without it.
            try:
                import pyzipper

                _zip_open = pyzipper.AESZipFile
            except ImportError:
                _zip_open = zipfile.ZipFile

            with _zip_open(full_path, "r") as zf:
                info_list = zf.infolist()
                needs_pw = any(i.flag_bits & 0x1 for i in info_list)

                if needs_pw:
                    if not password:
                        return jsonify({"error": "password_required"}), 401
                    pw_bytes = password.encode("utf-8")
                    zf.setpassword(pw_bytes)
                    # Verify the password by reading the first non-empty file.
                    # pyzipper raises RuntimeError("Bad password") on wrong password.
                    # stdlib raises RuntimeError("Bad password") for standard crypto
                    # or BadZipFile — but cannot verify AES (use pyzipper for that).
                    for info in info_list:
                        if not info.is_dir() and info.file_size > 0:
                            try:
                                zf.read(info.filename)
                            except RuntimeError as e:
                                if (
                                    "password" in str(e).lower()
                                    or "bad" in str(e).lower()
                                ):
                                    return jsonify({"error": "wrong_password"}), 401
                                raise
                            except zipfile.BadZipFile:
                                return jsonify({"error": "wrong_password"}), 401
                            break

                entries, total_size = [], 0
                for info in info_list[:MAX_ENTRIES]:
                    total_size += info.file_size
                    entries.append(
                        {
                            "name": html_lib.escape(info.filename),
                            "is_dir": info.is_dir(),
                            "size": info.file_size,
                            "compressed_size": info.compress_size,
                            "modified": _fmt_zip_date(info.date_time),
                        }
                    )

                return jsonify(
                    {
                        "type": "zip",
                        "entries": entries,
                        "total_entries": len(info_list),
                        "total_size": total_size,
                        "truncated": len(info_list) > MAX_ENTRIES,
                        "encrypted": needs_pw,
                    }
                )

        # ── TAR family (no native password support) ───────────────────────────
        elif ext in (
            "tar",
            "tar.gz",
            "tgz",
            "tar.bz2",
            "tbz2",
            "tar.xz",
            "txz",
            "gz",
            "bz2",
        ):
            import tarfile as _tf

            mode_map = {
                "tar": "r:",
                "tar.gz": "r:gz",
                "tgz": "r:gz",
                "tar.bz2": "r:bz2",
                "tbz2": "r:bz2",
                "tar.xz": "r:xz",
                "txz": "r:xz",
                "gz": "r:gz",
                "bz2": "r:bz2",
            }
            mode = mode_map.get(ext, "r:*")
            try:
                with _tf.open(full_path, mode) as tf:
                    members = tf.getmembers()
                    entries, total_size = [], 0
                    for m in members[:MAX_ENTRIES]:
                        is_dir = m.isdir()
                        size = 0 if is_dir else m.size
                        total_size += size
                        try:
                            modified = (
                                datetime.datetime.fromtimestamp(m.mtime).strftime(
                                    "%Y-%m-%d %H:%M"
                                )
                                if m.mtime
                                else ""
                            )
                        except Exception:
                            modified = ""
                        name = (
                            m.name
                            if not (is_dir and not m.name.endswith("/"))
                            else m.name + "/"
                        )
                        entries.append(
                            {
                                "name": html_lib.escape(name),
                                "is_dir": is_dir,
                                "size": size,
                                "compressed_size": None,
                                "modified": modified,
                            }
                        )
                    return jsonify(
                        {
                            "type": "tar",
                            "entries": entries,
                            "total_entries": len(members),
                            "total_size": total_size,
                            "truncated": len(members) > MAX_ENTRIES,
                            "encrypted": False,
                        }
                    )
            except _tf.TarError as e:
                return jsonify({"error": f"Could not read archive: {e}"}), 500

        # ── 7Z ────────────────────────────────────────────────────────────────
        elif ext == "7z":
            try:
                import py7zr
                import py7zr.exceptions as _7zexc
            except ImportError:
                return (
                    jsonify(
                        {
                            "error": "Missing package 'py7zr'. Install with: pip install py7zr"
                        }
                    ),
                    500,
                )

            # FIX: py7zr's needs_password() only returns True when the archive
            # was opened WITH a password — useless for probing. Instead, open
            # without a password and call list(); if the archive is encrypted,
            # py7zr raises PasswordRequired during list() / decompressor setup.
            needs_pw = False
            try:
                with py7zr.SevenZipFile(full_path, mode="r") as _probe:
                    _probe.list()  # raises PasswordRequired if encrypted
            except _7zexc.PasswordRequired:
                needs_pw = True
            except Exception:
                pass  # other errors caught later when opening for real

            if needs_pw and not password:
                return jsonify({"error": "password_required"}), 401

            try:
                kwargs = {"mode": "r"}
                if password:
                    kwargs["password"] = password

                with py7zr.SevenZipFile(full_path, **kwargs) as archive:
                    # list() will raise PasswordRequired if still wrong
                    try:
                        file_list = archive.list()
                    except _7zexc.PasswordRequired:
                        # Shouldn't happen for correct password, but be safe
                        if needs_pw:
                            return jsonify({"error": "wrong_password"}), 401
                        raise

                    entries, total_size = [], 0
                    for info in file_list[:MAX_ENTRIES]:
                        is_dir = info.is_directory
                        size = (info.uncompressed or 0) if not is_dir else 0
                        compressed = (info.compressed or 0) if not is_dir else 0
                        total_size += size
                        try:
                            modified = (
                                info.creationtime.strftime("%Y-%m-%d %H:%M")
                                if info.creationtime
                                else ""
                            )
                        except Exception:
                            modified = ""
                        name = info.filename
                        if is_dir and not name.endswith("/"):
                            name += "/"
                        entries.append(
                            {
                                "name": html_lib.escape(name),
                                "is_dir": is_dir,
                                "size": size,
                                "compressed_size": compressed,
                                "modified": modified,
                            }
                        )
                    return jsonify(
                        {
                            "type": "7z",
                            "entries": entries,
                            "total_entries": len(file_list),
                            "total_size": total_size,
                            "truncated": len(file_list) > MAX_ENTRIES,
                            "encrypted": needs_pw,
                        }
                    )
            except _7zexc.PasswordRequired:
                return jsonify({"error": "wrong_password"}), 401
            except (_7zexc.Bad7zFile, _7zexc.CrcError) as e:
                if needs_pw:
                    return jsonify({"error": "wrong_password"}), 401
                return jsonify({"error": f"Could not read archive: {e}"}), 500
            except Exception as e:
                err = str(e).lower()
                if needs_pw and any(
                    k in err for k in ("password", "wrong", "bad", "crc", "decrypt")
                ):
                    return jsonify({"error": "wrong_password"}), 401
                return jsonify({"error": f"Could not read archive: {e}"}), 500

        # ── RAR ───────────────────────────────────────────────────────────────
        elif ext == "rar":
            try:
                import rarfile
            except ImportError:
                return (
                    jsonify(
                        {
                            "error": "Missing package 'rarfile'. Install with: pip install rarfile"
                        }
                    ),
                    500,
                )

            try:
                rf = rarfile.RarFile(full_path)
            except rarfile.NotRarFile:
                return jsonify({"error": "Not a valid RAR file"}), 400
            except Exception as e:
                return jsonify({"error": f"Could not open RAR: {e}"}), 500

            with rf:
                parser = rf._file_parser
                header_encrypted = bool(parser and parser.has_header_encryption())

                # ── Case 1: Header encryption (WinRAR "Encrypt file names") ──
                # The file listing itself is encrypted. needs_password() returns
                # False here because no entries were parsed — check the parser flag.
                if header_encrypted:
                    if not password:
                        return jsonify({"error": "password_required"}), 401

                    # FIX: setpassword() takes a str, NOT bytes.
                    # For RAR5, setpassword() re-parses the archive and immediately
                    # raises RarWrongPassword if the PBKDF2 check fails.
                    try:
                        rf.setpassword(password)
                    except rarfile.RarWrongPassword:
                        return jsonify({"error": "wrong_password"}), 401
                    except rarfile.BadRarFile:
                        return jsonify({"error": "wrong_password"}), 401

                    # After correct setpassword, infolist() should now be populated
                    try:
                        info_list = rf.infolist()
                    except (rarfile.PasswordRequired, rarfile.BadRarFile):
                        return jsonify({"error": "wrong_password"}), 401

                    # Paranoia check: if the list is still empty the password was wrong
                    # (can happen with RAR3 header encryption where AES produces garbage)
                    if not info_list:
                        return jsonify({"error": "wrong_password"}), 401

                # ── Case 2: File-level encryption (content encrypted, list visible) ──
                else:
                    info_list = rf.infolist()
                    needs_pw = rf.needs_password()

                    if needs_pw:
                        if not password:
                            return jsonify({"error": "password_required"}), 401

                        # FIX: str, not bytes
                        rf.setpassword(password)

                        # Verify by attempting to read the first non-empty file
                        verified = False
                        for info in info_list:
                            if not info.is_dir() and info.file_size > 0:
                                try:
                                    rf.read(info.filename)
                                    verified = True
                                except (
                                    rarfile.RarWrongPassword,
                                    rarfile.RarCRCError,
                                    rarfile.BadRarFile,
                                    rarfile.PasswordRequired,
                                ):
                                    return jsonify({"error": "wrong_password"}), 401
                                break
                        # No extractable files to verify against — trust the listing
                        if not verified and needs_pw:
                            pass  # Accept — can't verify an all-directory archive

                # ── Build entry list ──────────────────────────────────────────
                entries, total_size = [], 0
                for info in info_list[:MAX_ENTRIES]:
                    is_dir = info.is_dir()
                    size = 0 if is_dir else (info.file_size or 0)
                    compressed = 0 if is_dir else (info.compress_size or 0)
                    total_size += size
                    try:
                        modified = (
                            info.date_time.strftime("%Y-%m-%d %H:%M")
                            if info.date_time
                            else ""
                        )
                    except Exception:
                        modified = ""
                    name = info.filename
                    if is_dir and not name.endswith("/"):
                        name += "/"
                    entries.append(
                        {
                            "name": html_lib.escape(name),
                            "is_dir": is_dir,
                            "size": size,
                            "compressed_size": compressed,
                            "modified": modified,
                        }
                    )

                return jsonify(
                    {
                        "type": "rar",
                        "entries": entries,
                        "total_entries": len(info_list),
                        "total_size": total_size,
                        "truncated": len(info_list) > MAX_ENTRIES,
                        "encrypted": header_encrypted
                        or (not header_encrypted and rf.needs_password()),
                    }
                )

        else:
            return jsonify({"error": f"Unsupported archive format: .{ext}"}), 400

    except Exception as e:
        logging.exception(f"Archive preview failed for {path}")
        return jsonify({"error": f"Preview failed: {e}"}), 500


@app.route("/bulk-download", methods=["POST"])
@login_required
def bulk_download():
    """Download multiple files and folders as a streaming ZIP file using zipstream-new"""
    try:
        print(f"📥 Bulk download request received from user: {current_user()}")

        # Handle both JSON and form data
        if request.is_json:
            data = request.get_json()
            print(f"📋 JSON Request data: {data}")

            if not data or "paths" not in data:
                print("❌ Error: No paths provided in JSON request")
                return jsonify({"error": "No paths provided"}), 400

            paths = data["paths"]
        else:
            # Handle form data
            print("📋 Form data request received")
            paths_json = request.form.get("paths")
            if not paths_json:
                print("❌ Error: No paths provided in form request")
                return jsonify({"error": "No paths provided"}), 400

            try:
                import json

                paths = json.loads(paths_json)
                print(f"📋 Form Request paths: {paths}")
            except json.JSONDecodeError:
                print("❌ Error: Invalid JSON in form paths")
                return jsonify({"error": "Invalid paths format"}), 400
        print(f"📁 Requested paths ({len(paths)} items): {paths}")

        if not paths:
            print("❌ Error: Empty paths list")
            return jsonify({"error": "Empty paths list"}), 400

        # Validate all paths
        print(f"🔍 Validating {len(paths)} paths...")
        invalid_paths = []
        valid_paths = []
        for path in paths:
            if not storage.is_safe_path(path):
                invalid_paths.append(path)
                print(f"⚠️  Invalid path detected: {path}")
            else:
                valid_paths.append(path)
                print(f"✅ Valid path: {path}")

        print(
            f"📊 Validation results: {len(valid_paths)} valid, {len(invalid_paths)} invalid"
        )

        if invalid_paths:
            return jsonify({"error": f"Invalid paths: {invalid_paths}"}), 400

        # Generate a filename for the ZIP based on selection
        if len(paths) == 1:
            # Single item - use its name
            base_name = os.path.basename(paths[0]) or "download"
        else:
            # Multiple items - use generic name with count
            base_name = f"bulk_download_{len(paths)}_items"

        zip_filename = f"{base_name}.zip"
        print(f"📦 Creating streaming ZIP file: {zip_filename}")

        # Capture session data before creating the generator (outside request context)
        session_id = session.get("session_id")
        if session_id:
            bulk_zip_progress[session_id] = {
                "current": 0,
                "total": len(paths),
                "done": False,
            }

        def generate_zip_stream():
            """Generator function to create ZIP file using zipstream-new for true streaming"""

            print(f"🗂️ Starting ZIP stream generation for {len(paths)} paths...")

            # Create zipstream object with optimized compression for large files
            zf = zipstream.ZipFile(
                mode="w", compression=zipstream.ZIP_DEFLATED, allowZip64=True
            )

            files_added = 0
            total_size = 0
            for i, path in enumerate(paths, 1):
                # Check for cancellation
                if session_id and bulk_zip_cancelled.get(session_id):
                    print(f"❌ ZIP generation cancelled for session {session_id}")
                    bulk_zip_cancelled.pop(session_id, None)
                    break

                print(f"📄 Processing item {i}/{len(paths)}: {path}")
                if session_id:
                    bulk_zip_progress[session_id]["current"] = i

                full_path = os.path.join(ROOT_DIR, path)
                if not os.path.exists(full_path):
                    print(f"⚠️  Path does not exist: {full_path}")
                    continue

                try:
                    if os.path.isfile(full_path):
                        # Add single file
                        arc_name = os.path.basename(full_path)
                        file_size = os.path.getsize(full_path)
                        total_size += file_size
                        print(
                            f"📄 Adding file to stream: {arc_name} ({file_size:,} bytes)"
                        )
                        zf.write(full_path, arcname=arc_name)
                        files_added += 1
                    elif os.path.isdir(full_path):
                        # Add directory recursively
                        dir_name = os.path.basename(full_path)
                        print(f"📁 Adding directory to stream: {dir_name}")
                        dir_files_added = 0

                        for root, dirs, files in os.walk(full_path):
                            # Calculate relative path for archive
                            rel_path = os.path.relpath(root, full_path)
                            if rel_path == ".":
                                arc_root = dir_name
                            else:
                                arc_root = os.path.join(dir_name, rel_path).replace(
                                    "\\", "/"
                                )

                            # Add all files in current directory
                            for file in files:
                                try:
                                    file_path = os.path.join(root, file)
                                    file_size = os.path.getsize(file_path)
                                    total_size += file_size
                                    arc_name = os.path.join(arc_root, file).replace(
                                        "\\", "/"
                                    )
                                    zf.write(file_path, arcname=arc_name)
                                    dir_files_added += 1
                                except (PermissionError, OSError) as e:
                                    print(f"⚠️  Skipped file {file_path}: {str(e)}")
                                    logging.warning(
                                        f"Skipped file {file_path}: {str(e)}"
                                    )
                                    continue

                            # Create empty directory entry if no files and no subdirs
                            if not files and not dirs:
                                zf.writestr(arc_root + "/", "")

                        print(f"📁 Directory added with {dir_files_added} files")
                        files_added += dir_files_added

                except (PermissionError, OSError) as e:
                    print(f"⚠️  Skipped item {full_path}: {str(e)}")
                    logging.warning(f"Skipped item {full_path}: {str(e)}")
                    continue

            print(
                f"✅ ZIP stream setup complete: {files_added} files queued for streaming"
            )
            if session_id:
                bulk_zip_progress[session_id]["done"] = True

            # Stream the ZIP file
            for chunk in zf:
                yield chunk

            print(f"� ZIP stream download completed")

        # Create response with streaming optimized for large files
        response = Response(
            generate_zip_stream(),
            mimetype="application/zip",
            headers={
                "Content-Disposition": f'attachment; filename="{zip_filename}"',
                "Cache-Control": "no-cache",
                "X-Accel-Buffering": "no",
                "Content-Encoding": "identity",
            },
        )

        print(f"🎉 Bulk download response ready for {len(paths)} items")
        logging.info(f"Bulk download initiated by {current_user()}: {len(paths)} items")
        logging.debug(f"Paths requested: {paths}")
        return response

    except Exception as e:
        print(f"❌ Bulk download error: {str(e)}")
        logging.error(f"Bulk download error: {str(e)}")
        return jsonify({"error": "Failed to create download"}), 500


@app.route("/upload", methods=["POST"])
@login_required
def upload():
    session_id = session.get("session_id")
    if not session_id:
        session_id = str(uuid.uuid4())
        session["session_id"] = session_id

    # Initialize these variables outside the try block to ensure they're available in the except blocks
    file_id = None
    filename = None

    try:
        role = get_role(current_user())
        if role != "readwrite":
            return "Permission denied", 403

        file_id = request.form.get("file_id")
        chunk_num = request.form.get("chunk_num")
        total_chunks = request.form.get("total_chunks")
        filename = request.form.get("filename", "")
        dest_path = request.form.get("dest_path", "")

        # Validate filename (must not be empty)
        if not filename:
            return "Filename is required", 400

        # Remove all sanitization, only check for empty and slashes
        if "/" in filename or "\\" in filename:
            return "Invalid filename", 400

        # Security check: ensure destination path is safe
        if dest_path and not storage.is_safe_path(dest_path):
            return "Invalid destination path", 400

        if (
            ENABLE_CHUNKED_UPLOADS
            and chunk_num is not None
            and total_chunks is not None
        ):
            # Chunked upload handling
            try:
                chunk_num = int(chunk_num)
                total_chunks = int(total_chunks)
            except ValueError:
                return "Invalid chunk parameters", 400

            if not file_id:
                return "File ID is required for chunked upload", 400

            # Track this upload
            chunk_tracker.track_upload(session_id, file_id)

            chunk = request.files.get("chunk")
            if not chunk:
                return "No chunk data received", 400

            chunk_data = chunk.read()
            if len(chunk_data) > CHUNK_SIZE:
                return f"Chunk too large (max {CHUNK_SIZE} bytes)", 413

            # Save chunk
            if not storage.save_chunk(file_id, chunk_num, chunk_data):
                # Cleanup on failure
                chunk_tracker.untrack_upload(session_id, file_id)
                storage.cleanup_chunks(file_id)
                return "Failed to save chunk", 500

            print(
                f"📦 Saved chunk {chunk_num + 1}/{total_chunks} for {filename} (ID: {file_id})"
            )

            # If this is the last chunk, queue for background assembly
            if chunk_num == total_chunks - 1:
                try:
                    # Save metadata for assembly worker
                    chunk_dir = os.path.join(ROOT_DIR, ".chunks", file_id)
                    metadata_file = os.path.join(chunk_dir, ".metadata")
                    metadata = {
                        "filename": filename,
                        "dest_path": dest_path,
                        "total_chunks": total_chunks,
                        "session_id": session_id,
                        "timestamp": time.time(),
                    }
                    with open(metadata_file, "w") as f:
                        json.dump(metadata, f)

                    # Add to background assembly queue
                    assembly_queue.add_job(
                        file_id, filename, dest_path, total_chunks, session_id
                    )

                    print(f"🔄 Queued {filename} for background assembly")
                    return (
                        jsonify(
                            {
                                "status": "upload_complete",
                                "message": f"Upload complete - processing {filename}...",
                                "file_id": file_id,
                                "assembly_queued": True,
                            }
                        ),
                        200,
                    )

                except Exception as e:
                    # Failed to queue assembly - cleanup
                    chunk_tracker.untrack_upload(session_id, file_id)
                    storage.cleanup_chunks(file_id)
                    print(f"❌ Failed to queue assembly for {filename}: {e}")
                    return f"Failed to queue file assembly: {str(e)}", 500

            return f"Chunk {chunk_num + 1}/{total_chunks} uploaded successfully", 200

        else:
            # Whole file upload handling
            uploaded_file = request.files.get("file")
            if not uploaded_file or uploaded_file.filename == "":
                return "No file selected", 400

            # Use provided filename or fall back to uploaded filename
            if not filename:
                filename = uploaded_file.filename
                if not filename:
                    return "Invalid filename", 400
            # Only check for slashes
            if "/" in filename or "\\" in filename:
                return "Invalid filename", 400

            # Construct target path
            target_dir = os.path.join(ROOT_DIR, dest_path) if dest_path else ROOT_DIR
            target_path = os.path.join(target_dir, filename)

            # Ensure target directory exists
            os.makedirs(target_dir, exist_ok=True)

            # Conflict check: return 409 if file exists and overwrite not explicitly requested
            overwrite = request.form.get("overwrite", "0")
            if (
                os.path.exists(target_path)
                and os.path.isfile(target_path)
                and overwrite != "1"
            ):
                return "File already exists", 409

            # Save file — retry on Windows file-lock errors (e.g. FastCopy holding a write lock).
            # Without this, save() blocks indefinitely waiting for the lock to release,
            # which stalls the Flask thread and freezes the entire upload queue.
            max_attempts = 3
            for attempt in range(max_attempts):
                try:
                    uploaded_file.stream.seek(0)
                    uploaded_file.save(target_path)
                    return "File uploaded successfully", 200
                except PermissionError as e:
                    if attempt < max_attempts - 1:
                        time.sleep(0.5)
                    else:
                        print(
                            f"❌ File locked after {max_attempts} attempts: {filename} — {e}"
                        )
                        return f"File is locked by another process: {str(e)}", 423
                except Exception as e:
                    print(f"❌ Failed to save whole file {filename}: {e}")
                    return f"Failed to save file: {str(e)}", 500

    except ClientDisconnected as e:
        print(
            f"👋 Client disconnected during upload of {filename or 'unknown file'} (ID: {file_id})"
        )
        # Untrack immediately (fast), but run the actual disk cleanup in the background
        # so this handler returns without blocking a Waitress thread on safe_rmtree.
        if file_id:
            chunk_tracker.untrack_upload(session_id, file_id)

            def _bg_disconnect_cleanup(fid):
                try:
                    storage.cleanup_chunks(fid)
                    print(f"🧹 Background disconnect cleanup done: {fid}")
                except Exception as ex:
                    print(f"⚠️ Background disconnect cleanup error for {fid}: {ex}")

            threading.Thread(
                target=_bg_disconnect_cleanup, args=(file_id,), daemon=True
            ).start()
        return "", 499

    except Exception as e:
        print(f"❌ Upload error: {e}")
        # If there was an error and we were tracking this upload, clean it up
        if file_id:
            chunk_tracker.untrack_upload(session_id, file_id)
            storage.cleanup_chunks(file_id)
        return f"Upload error: {str(e)}", 500


@app.route("/cleanup_chunks", methods=["POST"])
@login_required
def cleanup_chunks():
    """Clean up unfinished chunk files"""
    session_id = session.get("session_id")

    try:
        role = get_role(current_user())
        if role != "readwrite":
            return jsonify({"error": "Permission denied"}), 403

        data = request.get_json()
        if not data or "file_id" not in data:
            return jsonify({"error": "File ID is required"}), 400

        file_id = data["file_id"]

        # Untrack and cleanup
        chunk_tracker.untrack_upload(session_id, file_id)

        # Clean up chunks directory for this file_id
        chunks_dir = os.path.join(ROOT_DIR, ".chunks", file_id)
        if os.path.exists(chunks_dir):
            try:
                shutil.rmtree(chunks_dir)
                print(f"🧹 Manual cleanup completed for: {file_id}")

                # Try to remove parent chunks directory if empty
                parent_chunks_dir = os.path.join(ROOT_DIR, ".chunks")
                if os.path.exists(parent_chunks_dir) and not os.listdir(
                    parent_chunks_dir
                ):
                    os.rmdir(parent_chunks_dir)
                    print("🧹 Removed empty chunks directory")

                return (
                    jsonify(
                        {"success": True, "message": f"Cleaned up chunks for {file_id}"}
                    ),
                    200,
                )
            except Exception as e:
                print(f"❌ Failed to cleanup chunks for {file_id}: {e}")
                return jsonify({"error": f"Failed to cleanup chunks: {str(e)}"}), 500
        else:
            return jsonify({"success": True, "message": "No chunks to cleanup"}), 200

    except Exception as e:
        print(f"❌ Cleanup error: {e}")
        return jsonify({"error": f"Cleanup error: {str(e)}"}), 500


@app.route("/cancel_upload", methods=["POST"])
@login_required
def cancel_upload():
    """Cancel an ongoing upload and clean up its chunks"""
    session_id = session.get("session_id")

    try:
        role = get_role(current_user())
        if role != "readwrite":
            return jsonify({"error": "Permission denied"}), 403

        data = request.get_json()
        if not data or "file_id" not in data:
            return jsonify({"error": "File ID is required"}), 400

        file_id = data["file_id"]
        filename = data.get("filename", "Unknown file")

        print(f"🚫 Cancelling upload: {file_id} ({filename})")

        # Untrack the upload
        chunk_tracker.untrack_upload(session_id, file_id)

        # Clean up chunks directory for this file_id
        chunks_dir = os.path.join(ROOT_DIR, ".chunks", file_id)
        if os.path.exists(chunks_dir):
            try:
                # Run deletion in a daemon thread so this endpoint returns immediately
                # and does NOT block a Waitress thread (safe_rmtree can stall on Windows
                # file locks, which previously exhausted the thread pool when multiple
                # cancellations arrived at the same time).
                def _bg_cleanup(cdir, fid):
                    try:
                        storage.safe_rmtree(cdir)
                        print(f"🧹 Background cancelled-upload cleanup done: {fid}")
                        parent = os.path.join(ROOT_DIR, ".chunks")
                        if os.path.exists(parent):
                            try:
                                if not os.listdir(parent):
                                    os.rmdir(parent)
                            except OSError:
                                pass
                    except Exception as ex:
                        print(f"⚠️ Background cleanup error for {fid}: {ex}")

                threading.Thread(
                    target=_bg_cleanup, args=(chunks_dir, file_id), daemon=True
                ).start()
                print(f"🧹 Queued background cleanup for cancelled upload: {file_id}")

                return (
                    jsonify(
                        {
                            "success": True,
                            "message": f"Upload cancelled and cleanup queued for {filename}",
                            "file_id": file_id,
                        }
                    ),
                    200,
                )
            except Exception as e:
                print(f"❌ Failed to queue cleanup for cancelled upload {file_id}: {e}")
                return (
                    jsonify({"error": f"Failed to cleanup cancelled upload: {str(e)}"}),
                    500,
                )
        else:
            # Upload was cancelled before any chunks were created
            return (
                jsonify(
                    {
                        "success": True,
                        "message": f"Upload cancelled for {filename}",
                        "file_id": file_id,
                    }
                ),
                200,
            )

    except Exception as e:
        print(f"❌ Cancel upload error: {e}")
        return jsonify({"error": f"Cancel upload error: {str(e)}"}), 500


@app.route("/admin/rebuild_cache", methods=["POST"])
@login_required
def admin_rebuild_cache():
    """Delete storage_index.json and trigger a fresh full walk to rebuild it"""
    try:
        role = get_role(current_user())
        if role != "readwrite":
            return jsonify({"error": "Permission denied"}), 403

        from file_monitor import get_file_monitor, CACHE_FILE
        from file_index import file_index_manager, FILE_INDEX_PATH
        import os

        # Delete storage_index.json
        if os.path.exists(CACHE_FILE):
            os.remove(CACHE_FILE)
            print(f"🗑️ Cache file deleted: {CACHE_FILE}")
        else:
            print("ℹ️ No cache file found — nothing to delete")

        # Delete file_index.json
        if os.path.exists(FILE_INDEX_PATH):
            os.remove(FILE_INDEX_PATH)
            file_index_manager.clear()
            print(f"🗑️ File index deleted: {FILE_INDEX_PATH}")
        else:
            print("ℹ️ No file index found — nothing to delete")

        # Trigger a fresh full reconciliation walk to rebuild both, and
        # force-push SSE so the UI updates without a manual page refresh.
        monitor = get_file_monitor()
        print("🚶 Rebuilding cache from scratch...")
        monitor._reconcile()
        from realtime_stats import trigger_storage_update

        trigger_storage_update(None, monitor.get_current_snapshot())

        fi_stats = file_index_manager.get_stats()
        return (
            jsonify(
                {
                    "success": True,
                    "message": (
                        f"Cache cleared and rebuilt: {monitor._file_count:,} files, "
                        f"{monitor._dir_count:,} dirs, {len(monitor._dir_info):,} folders indexed. "
                        f'File index: {fi_stats["indexed_folders"]:,} large folder(s) indexed '
                        f'({fi_stats["total_entries"]:,} entries, threshold={fi_stats["threshold"]})'
                    ),
                }
            ),
            200,
        )

    except Exception as e:
        print(f"❌ Error during cache cleanup: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/cleanup_chunks", methods=["POST"])
@login_required
def admin_cleanup_chunks():
    """Admin endpoint to trigger comprehensive chunk cleanup"""
    from storage import manual_chunks_cleanup, emergency_cleanup_all

    try:
        role = get_role(current_user())
        if role != "readwrite":
            return jsonify({"error": "Permission denied"}), 403

        print("🧹 Starting comprehensive chunk cleanup...")

        # Get stats before cleanup
        stats_before = chunk_tracker.get_stats()

        # Cleanup orphaned chunks
        chunk_tracker.cleanup_orphaned_chunks()

        # Cleanup interrupted uploads
        chunk_tracker.cleanup_interrupted_uploads()

        # Get active assembly jobs to protect them from cleanup
        active_assembly_jobs = get_protected_files()

        if active_assembly_jobs:
            print(
                f"🔐 Manual cleanup protecting {len(active_assembly_jobs)} files currently being assembled"
            )

        # Cleanup old chunks (aggressive - 30 minutes)
        storage.cleanup_old_chunks(
            max_age_hours=0.5, protected_files=active_assembly_jobs
        )

        # Get stats after cleanup
        stats_after = chunk_tracker.get_stats()

        # Run enhanced manual cleanup
        manual_success = manual_chunks_cleanup()

        print(f"🧹 Comprehensive cleanup completed")
        print(
            f"   Sessions: {stats_before['active_sessions']} -> {stats_after['active_sessions']}"
        )
        print(
            f"   Uploads: {stats_before['active_uploads']} -> {stats_after['active_uploads']}"
        )

        return (
            jsonify(
                {
                    "success": True,
                    "message": (
                        "Comprehensive cleanup completed successfully"
                        if manual_success
                        else "Cleanup completed with some warnings"
                    ),
                    "stats_before": stats_before,
                    "stats_after": stats_after,
                    "manual_cleanup_success": manual_success,
                }
            ),
            200,
        )

    except Exception as e:
        print(f"❌ Error in comprehensive cleanup: {e}")
        # Try emergency cleanup as fallback
        try:
            emergency_cleanup_all()
            return (
                jsonify(
                    {
                        "success": True,
                        "message": f"Standard cleanup failed, emergency cleanup performed: {str(e)}",
                        "emergency_cleanup": True,
                    }
                ),
                200,
            )
        except Exception as emergency_error:
            return (
                jsonify(
                    {
                        "error": f"All cleanup methods failed: {str(e)} | Emergency: {str(emergency_error)}"
                    }
                ),
                500,
            )


@app.route("/admin/chunk_stats", methods=["GET"])
@login_required
def chunk_stats():
    """Get chunk tracking statistics"""
    try:
        role = get_role(current_user())
        if role != "readwrite":
            return jsonify({"error": "Permission denied"}), 403

        stats = chunk_tracker.get_stats()

        # Add filesystem stats
        chunks_dir = os.path.join(ROOT_DIR, ".chunks")
        filesystem_chunks = []
        if os.path.exists(chunks_dir):
            try:
                filesystem_chunks = [
                    d
                    for d in os.listdir(chunks_dir)
                    if os.path.isdir(os.path.join(chunks_dir, d))
                ]
            except OSError:
                pass

        stats["filesystem_chunks"] = len(filesystem_chunks)
        stats["chunk_directories"] = filesystem_chunks

        return jsonify(stats), 200

    except Exception as e:
        print(f"❌ Error getting chunk stats: {e}")
        return jsonify({"error": f"Stats error: {str(e)}"}), 500


@app.route("/admin/upload_status", methods=["GET"])
@login_required
def upload_status():
    """Get current upload status for UI updates"""
    try:
        role = get_role(current_user())

        # Allow readonly users to check auth status, but return limited info
        if role == "readonly":
            return jsonify(
                {
                    "authenticated": True,
                    "role": "readonly",
                    "has_active_uploads": False,
                    "session_has_active": False,
                    "total_active_sessions": 0,
                    "can_upload": False,
                }
            )

        if role != "readwrite":
            return jsonify({"error": "Permission denied"}), 403

        session_id = session.get("session_id")
        stats = chunk_tracker.get_stats()

        # Check if current session has active uploads
        session_has_active = False
        if session_id and session_id in chunk_tracker.active_uploads:
            session_has_active = len(chunk_tracker.active_uploads[session_id]) > 0

        return (
            jsonify(
                {
                    "has_active_uploads": stats["active_uploads"] > 0,
                    "session_has_active": session_has_active,
                    "total_active_sessions": stats["active_sessions"],
                    "total_active_uploads": stats["active_uploads"],
                }
            ),
            200,
        )

    except Exception as e:
        print(f"❌ Error getting upload status: {e}")
        return jsonify({"error": f"Status error: {str(e)}"}), 500


@app.route("/api/storage_stats", methods=["GET"])
@login_required
def storage_stats_api():
    """Get storage statistics - INSTANT VERSION using cached data"""
    try:
        print(
            f"📊 INSTANT Storage stats API called by user: {session.get('username', 'unknown')}"
        )

        # Use cached snapshot for instant response
        from file_monitor import get_file_monitor

        file_monitor = get_file_monitor()
        current_snapshot = file_monitor.get_current_snapshot()

        # Get fast disk stats only (no file counting)
        from realtime_stats import StorageStatsEventManager

        event_manager = StorageStatsEventManager()
        disk_stats = event_manager._get_fast_disk_stats()

        # Build instant stats response
        if current_snapshot:
            stats = {
                "total_space": disk_stats["total_space"],
                "used_space": disk_stats["used_space"],
                "free_space": disk_stats["free_space"],
                "file_count": current_snapshot.file_count,
                "dir_count": current_snapshot.dir_count,
                "content_size": current_snapshot.total_size,
            }
        else:
            # Fallback instant stats
            stats = {
                "total_space": disk_stats["total_space"],
                "used_space": disk_stats["used_space"],
                "free_space": disk_stats["free_space"],
                "file_count": 0,
                "dir_count": 0,
                "content_size": 0,
            }

        print(
            f"📊 INSTANT storage stats returned: files={stats['file_count']}, dirs={stats['dir_count']}"
        )
        return jsonify(stats), 200

    except Exception as e:
        print(f"❌ Error getting instant storage stats: {e}")
        import traceback

        traceback.print_exc()
        return jsonify({"error": f"Storage stats error: {str(e)}"}), 500


@app.route("/api/storage_stats_slow", methods=["GET"])
@login_required
def storage_stats_slow_api():
    """Get storage statistics - SLOW VERSION with full file counting"""
    try:
        print(
            f"📊 SLOW Storage stats API called by user: {session.get('username', 'unknown')}"
        )
        stats = storage.get_storage_stats()
        print(f"📊 SLOW storage stats calculated: {stats}")
        return jsonify(stats), 200

    except Exception as e:
        print(f"❌ Error getting slow storage stats: {e}")
        import traceback

        traceback.print_exc()
        return jsonify({"error": f"Storage stats error: {str(e)}"}), 500


@app.route("/api/storage_stats_debug", methods=["GET"])
def storage_stats_debug():
    """Debug version of storage stats without authentication"""
    try:
        print("🔧 Debug storage stats API called (no auth required)")
        stats = storage.get_storage_stats()
        print(f"🔧 Debug storage stats calculated: {stats}")
        return (
            jsonify(
                {
                    "debug": True,
                    "platform": os.name,
                    "has_statvfs": hasattr(os, "statvfs"),
                    "root_dir": storage.ROOT_DIR,
                    "stats": stats,
                }
            ),
            200,
        )

    except Exception as e:
        print(f"❌ Error in debug storage stats: {e}")
        import traceback

        traceback.print_exc()
        return jsonify({"error": f"Debug storage stats error: {str(e)}"}), 500


@app.route("/api/storage_stats_stream", methods=["GET"])
def storage_stats_stream():
    """Server-Sent Events endpoint for real-time storage stats"""
    if not is_logged_in():
        return jsonify({"error": "Authentication required"}), 401

    print(f"📡 SSE connection established for user: {current_user()}")
    return storage_stats_sse()


@app.route("/api/storage_stats_poll", methods=["GET"])
def storage_stats_poll():
    """Polling endpoint for storage stats - fallback when SSE fails"""
    if not is_logged_in():
        return jsonify({"error": "Authentication required"}), 401

    try:
        from file_monitor import get_file_monitor

        file_monitor = get_file_monitor()

        # Get current timestamp for comparison
        last_check = request.args.get("last_check", type=float, default=0)

        # For initial load (last_check=0), provide instant cached stats
        if last_check == 0:
            print("📊 Initial polling request - providing instant cached stats")
            current_time = time.time()

            # Get quick disk stats only
            from realtime_stats import StorageStatsEventManager

            event_manager = StorageStatsEventManager()
            disk_stats = event_manager._get_fast_disk_stats()

            # Use cached snapshot if available, otherwise provide placeholder
            current_snapshot = file_monitor.get_current_snapshot()
            if current_snapshot:
                file_count = current_snapshot.file_count
                dir_count = current_snapshot.dir_count
                total_size = current_snapshot.total_size
            else:
                # Provide instant placeholder stats
                file_count = 0
                dir_count = 0
                total_size = 0

            response_data = {
                "type": "polling_response",
                "timestamp": current_time,
                "changed": True,  # Always true for initial load
                "data": {
                    "file_count": file_count,
                    "dir_count": dir_count,
                    "total_size": total_size,
                    "content_size": total_size,
                    "last_modified": current_time,
                    "total_space": disk_stats["total_space"],
                    "free_space": disk_stats["free_space"],
                    "used_space": disk_stats["used_space"],
                    "changes": {
                        "files_changed": 0,
                        "dirs_changed": 0,
                        "size_changed": 0,
                        "content_changed": False,
                        "mtime_changed": False,
                    },
                },
            }

            print(f"📊 Instant polling response: files={file_count}, dirs={dir_count}")
            return jsonify(response_data), 200

        # Regular polling check for changes
        current_snapshot = file_monitor.get_current_snapshot()
        current_time = time.time()

        # Always return current stats, but include a 'changed' flag
        has_changes = False
        changes_data = {"files_changed": 0, "dirs_changed": 0, "size_changed": 0}

        if current_snapshot and current_snapshot.timestamp > last_check:
            has_changes = True

            # Get the last known file/dir counts from the polling history
            # Use a simple session-based tracking to reduce false positives
            last_known_files = request.args.get("last_files", type=int, default=0)
            last_known_dirs = request.args.get("last_dirs", type=int, default=0)

            # Calculate actual count changes
            files_diff = (
                current_snapshot.file_count - last_known_files
                if last_known_files > 0
                else 0
            )
            dirs_diff = (
                current_snapshot.dir_count - last_known_dirs
                if last_known_dirs > 0
                else 0
            )

            # Only report specific changes if we have meaningful differences
            if abs(files_diff) > 0 or abs(dirs_diff) > 0:
                # Real file/folder count change detected
                changes_data = {
                    "files_changed": files_diff,
                    "dirs_changed": dirs_diff,
                    "size_changed": 0,  # Size changes are complex to calculate
                    "content_changed": True,
                    "mtime_changed": True,
                }
            else:
                # Timestamp changed but no count changes - likely system noise
                # Report as minor content change without specific counts
                changes_data = {
                    "files_changed": 0,  # No count change
                    "dirs_changed": 0,  # No count change
                    "size_changed": 0,  # No size change claimed
                    "content_changed": True,  # Something changed (timestamp)
                    "mtime_changed": True,  # Modification time changed
                }

        # Debug logging for timestamp comparison
        print(
            f"📊 Polling debug: last_check={last_check}, snapshot_timestamp={current_snapshot.timestamp if current_snapshot else 'None'}, has_changes={has_changes}"
        )

        # Get disk stats
        from realtime_stats import StorageStatsEventManager

        event_manager = StorageStatsEventManager()
        disk_stats = event_manager._get_fast_disk_stats()

        response_data = {
            "type": "polling_response",
            "timestamp": current_time,
            "changed": has_changes,
            "data": {
                "file_count": current_snapshot.file_count if current_snapshot else 0,
                "dir_count": current_snapshot.dir_count if current_snapshot else 0,
                "total_size": current_snapshot.total_size if current_snapshot else 0,
                "content_size": current_snapshot.total_size if current_snapshot else 0,
                "last_modified": (
                    current_snapshot.last_modified if current_snapshot else current_time
                ),
                "total_space": disk_stats["total_space"],
                "free_space": disk_stats["free_space"],
                "used_space": disk_stats["used_space"],
                "changes": changes_data,  # Add changes field for frontend
            },
        }

        print(
            f"📊 Polling response: changed={has_changes}, files={response_data['data']['file_count']}"
        )
        return jsonify(response_data), 200

    except Exception as e:
        print(f"❌ Error in polling endpoint: {e}")
        return jsonify({"error": f"Polling error: {str(e)}"}), 500


@app.route("/api/monitoring_status", methods=["GET"])
def monitoring_status():
    """Get current monitoring system status"""
    if not is_logged_in():
        return jsonify({"error": "Authentication required"}), 401

    try:
        event_manager = get_event_manager()
        return (
            jsonify(
                {
                    "monitoring_active": file_monitor.monitoring,
                    "connected_clients": event_manager.get_client_count(),
                    "last_check": getattr(file_monitor, "last_check_time", None),
                    "total_checks": getattr(file_monitor, "check_count", 0),
                }
            ),
            200,
        )
    except Exception as e:
        print(f"❌ Error getting monitoring status: {e}")
        return jsonify({"error": f"Monitoring status error: {str(e)}"}), 500


@app.route("/api/disk_stats_fast", methods=["GET"])
def disk_stats_fast():
    """Fast disk stats only (no file counting) - no auth required"""
    try:
        print("📊 Fast disk stats request")

        # Get only disk usage stats, skip file counting
        disk_usage_path = storage.ROOT_DIR

        # Special handling for Android/Termux
        if "TERMUX_VERSION" in os.environ or os.path.exists("/data/data/com.termux"):
            android_storage_paths = [
                "/storage/emulated/0",
                "/sdcard",
                "/storage/self/primary",
            ]

            for path in android_storage_paths:
                if os.path.exists(path) and os.access(path, os.R_OK):
                    disk_usage_path = path
                    break

        # Get disk usage only
        if hasattr(os, "statvfs"):  # Unix-like systems
            try:
                stat = os.statvfs(disk_usage_path)
                total = stat.f_blocks * stat.f_frsize
                free = stat.f_bavail * stat.f_frsize
                used = total - free
            except OSError:
                # Fallback to shutil
                import shutil

                total, used, free = shutil.disk_usage(disk_usage_path)
        else:  # Windows
            import shutil

            total, used, free = shutil.disk_usage(storage.ROOT_DIR)

        return (
            jsonify(
                {
                    "total_space": total,
                    "used_space": used,
                    "free_space": free,
                    "file_count": "counting...",  # Will be updated by full stats
                    "dir_count": "counting...",
                    "content_size": "counting...",
                }
            ),
            200,
        )

    except Exception as e:
        print(f"❌ Error in fast disk stats: {e}")
        return jsonify({"error": f"Fast disk stats error: {str(e)}"}), 500


@app.route("/api/health_check", methods=["GET"])
def health_check():
    """Simple health check endpoint that doesn't require authentication"""
    return (
        jsonify(
            {
                "status": "ok",
                "platform": os.name,
                "has_statvfs": hasattr(os, "statvfs"),
                "root_dir": ROOT_DIR,
                "timestamp": time.time(),
            }
        ),
        200,
    )


@app.route("/api/search", methods=["GET"])
@login_required
def search_files():
    """Deep search through all folders for files/folders matching query"""
    query = request.args.get("q", "").strip()
    if not query:
        return jsonify({"results": [], "query": query}), 200

    try:
        query_lower = query.lower()
        results = []
        search_start = time.time()
        max_results = 100  # Limit results to avoid overwhelming UI

        # Walk through all directories starting from ROOT_DIR
        for root, dirs, files in os.walk(ROOT_DIR):
            # Stop if we've found enough results
            if len(results) >= max_results:
                break

            # Calculate relative path from ROOT_DIR
            rel_path = os.path.relpath(root, ROOT_DIR)
            if rel_path == ".":
                rel_path = ""

            # Search in folder names
            for dirname in dirs[:]:  # Use slice to allow modification during iteration
                if query_lower in dirname.lower():
                    folder_path = (
                        os.path.join(rel_path, dirname) if rel_path else dirname
                    )
                    try:
                        full_path = os.path.join(root, dirname)
                        stat_info = os.stat(full_path)

                        results.append(
                            {
                                "name": dirname,
                                "path": folder_path.replace("\\", "/"),
                                "type": "folder",
                                "is_dir": True,
                                "size": 0,
                                "modified": datetime.fromtimestamp(
                                    stat_info.st_mtime
                                ).strftime("%Y-%m-%d %H:%M:%S"),
                                "match_type": "name",
                            }
                        )
                    except (OSError, IOError):
                        continue  # Skip inaccessible folders

                    if len(results) >= max_results:
                        break

            # Search in file names
            for filename in files:
                if len(results) >= max_results:
                    break

                if query_lower in filename.lower():
                    file_path = (
                        os.path.join(rel_path, filename) if rel_path else filename
                    )
                    try:
                        full_path = os.path.join(root, filename)
                        stat_info = os.stat(full_path)

                        # Get file extension for type
                        _, ext = os.path.splitext(filename)
                        file_type = ext[1:].upper() if ext else "FILE"

                        results.append(
                            {
                                "name": filename,
                                "path": file_path.replace("\\", "/"),
                                "type": file_type,
                                "is_dir": False,
                                "size": stat_info.st_size,
                                "modified": datetime.fromtimestamp(
                                    stat_info.st_mtime
                                ).strftime("%Y-%m-%d %H:%M:%S"),
                                "match_type": "name",
                            }
                        )
                    except (OSError, IOError):
                        continue  # Skip inaccessible files

        search_time = time.time() - search_start

        return (
            jsonify(
                {
                    "results": results,
                    "query": query,
                    "total_found": len(results),
                    "search_time": round(search_time, 3),
                    "truncated": len(results) >= max_results,
                }
            ),
            200,
        )

    except Exception as e:
        print(f"❌ Search error: {str(e)}")
        return jsonify({"error": f"Search failed: {str(e)}"}), 500


@app.route("/api/dir_info/", defaults={"path": ""})
@app.route("/api/dir_info/<path:path>")
@login_required
def dir_info(path):
    """
    Returns folder size and item count.
    Hits the in-memory index instantly if indexed.
    Falls back to live walk for brand-new folders not yet in the index,
    then stores the result back so subsequent requests are instant.
    """
    if path and not storage.is_safe_path(path):
        return jsonify({"error": "Invalid path"}), 400
    try:
        info = storage.get_dir_info(path)

        # If this was a live walk fallback, store it back into the monitor index
        # so the next request for this path is instant
        try:
            from file_monitor import get_file_monitor

            monitor = get_file_monitor()
            rel_path = path.replace("\\", "/").strip("/")
            if monitor.get_dir_info(rel_path) is None:
                with monitor.lock:
                    monitor._dir_info[rel_path] = {
                        "file_count": info["file_count"],
                        "dir_count": info["dir_count"],
                        "total_size": info["total_size"],
                    }
                print(f"📥 Stored live walk result for '{rel_path}' into index")
        except Exception:
            pass

        return jsonify(info), 200
    except Exception as e:
        print(f"❌ Error getting dir info for {path}: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/assembly_status", methods=["GET"])
@login_required
def get_assembly_status():
    """Get all assembly jobs for current session"""
    session_id = session.get("session_id")
    if not session_id:
        return jsonify({"jobs": []}), 200

    jobs = assembly_queue.get_jobs_for_session(session_id)
    job_data = []

    for job in jobs:
        job_data.append(
            {
                "file_id": job.file_id,
                "filename": job.filename,
                "status": job.status,
                "created_at": job.created_at,
                "error_message": job.error_message,
            }
        )

    return jsonify({"jobs": job_data}), 200


@app.route("/api/protect_assembly/<file_id>", methods=["POST"])
@login_required
def protect_assembly_job(file_id):
    """Mark an assembly job as protected from cleanup"""
    session_id = session.get("session_id")
    if not session_id:
        return jsonify({"error": "No session ID"}), 400

    # Check if this job belongs to the current session
    job = assembly_queue.get_job_status(file_id)
    if job and job.session_id == session_id:
        # Re-track this upload to prevent cleanup
        chunk_tracker.track_upload(session_id, file_id)
        print(f"🔐 Protected assembly job {file_id} from cleanup")
        return jsonify({"status": "protected"}), 200

    return jsonify({"error": "Job not found or access denied"}), 404


@app.route("/api/assembly_status/<file_id>", methods=["GET"])
@login_required
def get_single_assembly_status(file_id):
    """Get status of a specific assembly job"""
    job = assembly_queue.get_job_status(file_id)

    if not job:
        return jsonify({"error": "Job not found"}), 404

    # Check if user owns this job
    session_id = session.get("session_id")
    if job.session_id != session_id:
        return jsonify({"error": "Access denied"}), 403

    return (
        jsonify(
            {
                "file_id": job.file_id,
                "filename": job.filename,
                "status": job.status,
                "created_at": job.created_at,
                "error_message": job.error_message,
            }
        ),
        200,
    )


@app.route("/api/files/", defaults={"path": ""})
@app.route("/api/files/<path:path>")
@login_required
def api_files(path):
    """API endpoint to get file listings as JSON"""
    try:
        # Comprehensive path validation: safety and existence
        if path and not storage.is_valid_path(path):
            return jsonify({"error": "Invalid path"}), 400

        role = get_role(current_user())
        items = storage.list_dir(path)

        response_data = {
            "success": True,
            "files": items,
            "current_path": path,
            "role": role,
        }

        return jsonify(response_data), 200

    except Exception as e:
        print(f"❌ Error in api_files: {e}")
        return jsonify({"error": "Failed to load files"}), 500


@app.route("/bulk_move", methods=["POST"])
@login_required
def bulk_move():
    """Move multiple files/folders to a new location"""
    try:
        role = get_role(current_user())
        if role != "readwrite":
            return jsonify({"error": "Permission denied"}), 403

        data = request.get_json()
        if not data or "paths" not in data:
            return jsonify({"error": "Paths are required"}), 400

        paths = data["paths"]
        destination = data.get("destination", "").strip()
        current_path = data.get("current_path", "")

        if not paths:
            return jsonify({"error": "No paths provided"}), 400

        # Validate destination path
        if destination and not storage.is_safe_path(destination):
            return jsonify({"error": "Invalid destination path"}), 400

        # conflict_resolutions maps filename -> 'overwrite' | 'rename' | 'skip'
        conflict_resolutions = data.get("conflict_resolutions", {})

        moved_count = 0
        errors = []

        def _find_free_name(dest_dir, filename):
            base, ext = os.path.splitext(filename)
            for i in range(1, 1000):
                candidate = f"{base} ({i}){ext}"
                if not os.path.exists(os.path.join(dest_dir, candidate)):
                    return candidate
            return f"{base} ({int(time.time())}){ext}"

        for source_path in paths:
            try:
                # Security check
                if not storage.is_safe_path(source_path):
                    errors.append(f"Invalid source path: {source_path}")
                    continue

                source_full = os.path.join(ROOT_DIR, source_path)
                if not os.path.exists(source_full):
                    errors.append(f"Source not found: {source_path}")
                    continue

                # Determine destination
                filename = os.path.basename(source_path)
                dest_dir = (
                    os.path.join(ROOT_DIR, destination) if destination else ROOT_DIR
                )
                dest_full = os.path.join(dest_dir, filename)

                # Create destination directory if it doesn't exist
                os.makedirs(dest_dir, exist_ok=True)

                # Handle conflict
                if os.path.exists(dest_full):
                    resolution = conflict_resolutions.get(filename, "error")
                    if resolution == "skip":
                        continue
                    elif resolution == "overwrite":
                        if os.path.isdir(dest_full):
                            shutil.rmtree(dest_full)
                        else:
                            os.remove(dest_full)
                    elif resolution == "rename":
                        dest_full = os.path.join(
                            dest_dir, _find_free_name(dest_dir, filename)
                        )
                    else:
                        errors.append(
                            f"Destination already exists: {os.path.join(destination, filename) if destination else filename}"
                        )
                        continue

                # Perform the move
                shutil.move(source_full, dest_full)
                moved_count += 1

            except Exception as e:
                errors.append(f"Failed to move {source_path}: {str(e)}")

        if errors:
            return (
                jsonify(
                    {
                        "moved_count": moved_count,
                        "errors": errors,
                        "error": f"Some items could not be moved. Moved {moved_count} items with {len(errors)} errors.",
                    }
                ),
                207,
            )  # Multi-status
        else:
            return jsonify({"moved_count": moved_count, "success": True}), 200

    except Exception as e:
        return jsonify({"error": f"Bulk move error: {str(e)}"}), 500


@app.route("/bulk_copy", methods=["POST"])
@login_required
def bulk_copy():
    """Copy multiple files/folders to a new location"""
    try:
        role = get_role(current_user())
        if role != "readwrite":
            return jsonify({"error": "Permission denied"}), 403

        data = request.get_json()
        if not data or "paths" not in data:
            return jsonify({"error": "Paths are required"}), 400

        paths = data["paths"]
        destination = data.get("destination", "").strip()
        current_path = data.get("current_path", "")

        if not paths:
            return jsonify({"error": "No paths provided"}), 400

        # Validate destination path
        if destination and not storage.is_safe_path(destination):
            return jsonify({"error": "Invalid destination path"}), 400

        # conflict_resolutions maps filename -> 'overwrite' | 'rename' | 'skip'
        conflict_resolutions = data.get("conflict_resolutions", {})

        copied_count = 0
        errors = []

        def _find_free_name_copy(dest_dir, filename):
            base, ext = os.path.splitext(filename)
            for i in range(1, 1000):
                candidate = f"{base} ({i}){ext}"
                if not os.path.exists(os.path.join(dest_dir, candidate)):
                    return candidate
            return f"{base} ({int(time.time())}){ext}"

        for source_path in paths:
            try:
                # Security check
                if not storage.is_safe_path(source_path):
                    errors.append(f"Invalid source path: {source_path}")
                    continue

                source_full = os.path.join(ROOT_DIR, source_path)
                if not os.path.exists(source_full):
                    errors.append(f"Source not found: {source_path}")
                    continue

                # Determine destination
                filename = os.path.basename(source_path)
                dest_dir = (
                    os.path.join(ROOT_DIR, destination) if destination else ROOT_DIR
                )
                dest_full = os.path.join(dest_dir, filename)

                # Create destination directory if it doesn't exist
                os.makedirs(dest_dir, exist_ok=True)

                # Handle conflict
                if os.path.exists(dest_full):
                    resolution = conflict_resolutions.get(
                        filename, "rename"
                    )  # default: auto-rename
                    if resolution == "skip":
                        continue
                    elif resolution == "overwrite":
                        if os.path.isdir(dest_full):
                            shutil.rmtree(dest_full)
                        else:
                            os.remove(dest_full)
                    else:  # 'rename' or default
                        dest_full = os.path.join(
                            dest_dir, _find_free_name_copy(dest_dir, filename)
                        )

                # Perform the copy
                if os.path.isdir(source_full):
                    shutil.copytree(source_full, dest_full)
                else:
                    shutil.copy2(source_full, dest_full)

                copied_count += 1

            except Exception as e:
                errors.append(f"Failed to copy {source_path}: {str(e)}")

        if copied_count > 0:
            _trigger_reconcile(settle=True)  # copytree fires a backlog storm

        if errors:
            return (
                jsonify(
                    {
                        "copied_count": copied_count,
                        "errors": errors,
                        "error": f"Some items could not be copied. Copied {copied_count} items with {len(errors)} errors.",
                    }
                ),
                207,
            )  # Multi-status
        else:
            return jsonify({"copied_count": copied_count, "success": True}), 200

    except Exception as e:
        return jsonify({"error": f"Bulk copy error: {str(e)}"}), 500


@app.route("/bulk_delete", methods=["POST"])
@login_required
def bulk_delete():
    """Delete multiple files/folders"""
    try:
        role = get_role(current_user())
        if role != "readwrite":
            return jsonify({"error": "Permission denied"}), 403

        data = request.get_json()
        if not data or "paths" not in data:
            return jsonify({"error": "Paths are required"}), 400

        paths = data["paths"]

        if not paths:
            return jsonify({"error": "No paths provided"}), 400

        deleted_count = 0
        errors = []

        for target_path in paths:
            try:
                # Security check
                if not storage.is_safe_path(target_path):
                    errors.append(f"Invalid path: {target_path}")
                    continue

                full_path = os.path.join(ROOT_DIR, target_path)
                if not os.path.exists(full_path):
                    errors.append(f"Path not found: {target_path}")
                    continue

                # Perform the deletion
                if os.path.isdir(full_path):
                    shutil.rmtree(full_path)
                else:
                    os.remove(full_path)

                deleted_count += 1

            except Exception as e:
                errors.append(f"Failed to delete {target_path}: {str(e)}")

        # Reconcile immediately so file/dir counts are corrected without waiting 15 min
        if deleted_count > 0:
            _trigger_reconcile()

        if errors:
            return (
                jsonify(
                    {
                        "deleted_count": deleted_count,
                        "errors": errors,
                        "error": f"Some items could not be deleted. Deleted {deleted_count} items with {len(errors)} errors.",
                    }
                ),
                207,
            )  # Multi-status
        else:
            return jsonify({"deleted_count": deleted_count, "success": True}), 200

    except Exception as e:
        return jsonify({"error": f"Bulk delete error: {str(e)}"}), 500


@app.route("/rename", methods=["POST"])
@login_required
def rename_item():
    """Rename a single file or folder"""
    try:
        role = get_role(current_user())
        if role != "readwrite":
            return jsonify({"error": "Permission denied"}), 403

        data = request.get_json()
        if not data or "old_path" not in data or "new_name" not in data:
            return jsonify({"error": "Old path and new name are required"}), 400

        old_path = data["old_path"]
        new_name = data["new_name"].strip()

        # Validate inputs
        if not new_name:
            return jsonify({"error": "New name cannot be empty"}), 400

        # Security checks
        if not storage.is_safe_path(old_path):
            return jsonify({"error": "Invalid old path"}), 400

        # Validate new name doesn't contain path separators or invalid characters
        if (
            "/" in new_name
            or "\\" in new_name
            or any(char in new_name for char in '<>:"|?*')
        ):
            return jsonify({"error": "Invalid characters in new name"}), 400

        # Check if old path exists
        old_full_path = os.path.join(ROOT_DIR, old_path)
        if not os.path.exists(old_full_path):
            return jsonify({"error": "Item not found"}), 404

        # Get the directory of the old path
        parent_dir = os.path.dirname(old_path)

        # Create new path
        new_path = os.path.join(parent_dir, new_name) if parent_dir else new_name
        new_full_path = os.path.join(ROOT_DIR, new_path)

        # Check if destination already exists
        if os.path.exists(new_full_path):
            return jsonify({"error": "An item with that name already exists"}), 409

        # Perform the rename
        try:
            os.rename(old_full_path, new_full_path)
            return (
                jsonify(
                    {
                        "success": True,
                        "message": f'Successfully renamed to "{new_name}"',
                        "old_path": old_path,
                        "new_path": new_path,
                        "new_name": new_name,
                    }
                ),
                200,
            )
        except OSError as e:
            return jsonify({"error": f"Failed to rename: {str(e)}"}), 500

    except Exception as e:
        return jsonify({"error": f"Rename error: {str(e)}"}), 500


@app.route("/mkdir", methods=["POST"])
@login_required
def mkdir():
    try:
        role = get_role(current_user())
        if role != "readwrite":
            return jsonify({"error": "Permission denied"}), 403

        foldername = request.form.get("foldername", "").strip()
        path = request.form.get("path", "")

        if not foldername:
            return jsonify({"error": "Folder name required"}), 400

        # Only replace slashes, preserve all other characters (including +, spaces, etc.)
        foldername = foldername.replace("/", "_").replace("\\", "_")
        if not foldername:
            return jsonify({"error": "Invalid folder name"}), 400

        # Security check: ensure path is safe
        if path and not storage.is_safe_path(path):
            return jsonify({"error": "Invalid path"}), 400

        created = storage.create_folder(path, foldername)
        if not created:
            return (
                jsonify({"error": "Folder already exists or could not be created"}),
                409,
            )
        else:
            return (
                jsonify(
                    {
                        "success": True,
                        "message": f'Folder "{foldername}" created successfully',
                    }
                ),
                200,
            )

    except Exception as e:
        return jsonify({"error": f"Error creating folder: {str(e)}"}), 500


@app.route("/delete", methods=["POST"])
@login_required
def delete():
    try:
        role = get_role(current_user())
        if role != "readwrite":
            flash("Permission denied")
            return redirect(url_for("index"))

        target_path = request.form.get("target_path")
        if not target_path:
            flash("Target path is required")
            return redirect(url_for("index"))

        # Security check: ensure path is safe
        if not storage.is_safe_path(target_path):
            flash("Invalid target path")
            return redirect(url_for("index"))

        if storage.delete_path(target_path):
            flash("Item deleted successfully")
            _trigger_reconcile()
        else:
            flash("Error deleting item")

        # Redirect to parent directory
        parent_path = "/".join(target_path.split("/")[:-1])
        return redirect(url_for("index", path=parent_path))

    except Exception as e:
        flash(f"Error deleting item: {str(e)}")
        return redirect(url_for("index"))


@app.route("/api/check_conflicts", methods=["POST"])
@login_required
def api_check_conflicts():
    """Check which of the given paths would conflict at the destination."""
    try:
        data = request.get_json()
        if not data:
            return jsonify({"error": "JSON body required"}), 400
        paths = data.get("paths", [])
        destination = data.get("destination", "").strip()

        if destination and not storage.is_safe_path(destination):
            return jsonify({"error": "Invalid destination path"}), 400

        conflicts = []
        for source_path in paths:
            if not storage.is_safe_path(source_path):
                continue
            filename = os.path.basename(source_path)
            dest_full = (
                os.path.join(ROOT_DIR, destination, filename)
                if destination
                else os.path.join(ROOT_DIR, filename)
            )
            if os.path.exists(dest_full):
                conflicts.append(
                    {
                        "source": source_path,
                        "name": filename,
                        "is_dir": os.path.isdir(dest_full),
                    }
                )
        return jsonify({"conflicts": conflicts})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/exists", methods=["GET"])
@login_required
def api_exists():
    """Check whether a path (file or folder) exists under ROOT_DIR."""
    path = request.args.get("path", "")
    if path and not storage.is_safe_path(path):
        return jsonify({"error": "Invalid path"}), 400
    full_path = os.path.join(ROOT_DIR, path) if path else ROOT_DIR
    exists = os.path.exists(full_path)
    is_dir = os.path.isdir(full_path) if exists else False
    return jsonify({"exists": exists, "is_dir": is_dir, "path": path})


@app.route("/api/speedtest/ping", methods=["GET"])
def speedtest_ping():
    # Just return OK for latency test
    return jsonify({"ok": True})


@app.route("/api/speedtest/upload", methods=["POST"])
def speedtest_upload():
    # Receive 25MiB data, measure time server-side if needed
    file = request.files.get("data")
    if not file:
        return jsonify({"error": "No data"}), 400
    # Optionally read to memory to simulate disk write
    file.read()
    return jsonify({"ok": True})


@app.route("/api/speedtest/download", methods=["GET"])
def speedtest_download():
    # Send 5MiB of zero bytes
    size = 5 * 1024 * 1024
    buf = io.BytesIO(b"\x00" * size)
    return send_file(
        buf,
        mimetype="application/octet-stream",
        as_attachment=True,
        download_name="speedtest.bin",
    )


@app.errorhandler(413)
def too_large(e):
    return "File too large", 413


@app.errorhandler(404)
def not_found(e):
    return render_template("404.html"), 404


@app.errorhandler(500)
def internal_error(e):
    return "Internal server error", 500


# Enhanced cleanup scheduler functions
def start_enhanced_cleanup_scheduler():
    """Start enhanced background thread for chunk cleanup"""

    def cleanup_worker():
        while True:
            try:
                # More frequent cleanup - every 15 minutes for stale chunks
                time.sleep(900)  # 15 minutes
                print("🧹 Running enhanced chunk cleanup...")

                # Get active assembly jobs to protect them from cleanup
                active_assembly_jobs = get_protected_files()

                if active_assembly_jobs:
                    print(
                        f"🔐 Protecting {len(active_assembly_jobs)} files from periodic cleanup"
                    )

                storage.cleanup_old_chunks(
                    max_age_hours=1, protected_files=active_assembly_jobs
                )  # Clean 1+ hour old chunks

                # Every 4th run (1 hour), do the full 24-hour cleanup
                cleanup_counter = getattr(cleanup_worker, "counter", 0) + 1
                cleanup_worker.counter = cleanup_counter

                if cleanup_counter % 4 == 0:  # Every hour
                    print("🧹 Running full chunk cleanup...")
                    storage.cleanup_old_chunks(
                        max_age_hours=24, protected_files=active_assembly_jobs
                    )

            except Exception as e:
                print(f"❌ Error in enhanced cleanup worker: {e}")

    cleanup_thread = threading.Thread(target=cleanup_worker, daemon=True)
    cleanup_thread.start()
    print("🧹 Started enhanced chunk cleanup scheduler (every 15 minutes)")


def start_orphan_cleanup_scheduler():
    """Start a background thread to cleanup orphaned chunks and detect interruptions"""

    def orphan_cleanup_worker():
        while True:
            try:
                time.sleep(300)  # Every 5 minutes
                chunk_tracker.cleanup_orphaned_chunks()
                chunk_tracker.cleanup_interrupted_uploads()
            except Exception as e:
                print(f"❌ Error in orphan cleanup worker: {e}")

    cleanup_thread = threading.Thread(target=orphan_cleanup_worker, daemon=True)
    cleanup_thread.start()
    print("🗑️ Started enhanced orphaned chunk cleanup scheduler (every 5 minutes)")


def assembly_worker():
    """Background worker that processes assembly jobs"""
    print("🔄 Assembly worker started")

    while True:
        try:
            # Get next job from queue (blocks until available)
            job = assembly_queue.job_queue.get(timeout=10)

            print(f"🔨 Processing assembly job: {job.filename} (ID: {job.file_id})")

            # Update job status to processing
            with assembly_queue.lock:
                if job.file_id in assembly_queue.active_jobs:
                    assembly_queue.active_jobs[job.file_id].status = "processing"

            try:
                # Perform the actual assembly
                success = storage.assemble_chunks(
                    job.file_id, job.filename, job.dest_path
                )

                if success:
                    assembly_queue.complete_job(job.file_id, success=True)
                    print(f"✅ Successfully assembled: {job.filename}")
                else:
                    assembly_queue.complete_job(
                        job.file_id,
                        success=False,
                        error_message="Assembly failed - see server logs",
                    )
                    print(f"❌ Assembly failed: {job.filename}")

            except Exception as e:
                error_msg = str(e)
                assembly_queue.complete_job(
                    job.file_id, success=False, error_message=error_msg
                )
                print(f"❌ Assembly error for {job.filename}: {error_msg}")

            # Mark queue task as done
            assembly_queue.job_queue.task_done()

            # When the assembly queue drains to zero, trigger an immediate reconcile
            # so the file count corrects itself right away. Use _trigger_reconcile (not
            # reconcile_async) so the SSE force-push fires even if watchdog already
            # updated the counters and _reconcile sees no drift.
            if assembly_queue.job_queue.empty() and not assembly_queue.active_jobs:
                _trigger_reconcile()

        except queue.Empty:
            # Timeout - cleanup old jobs periodically
            assembly_queue.cleanup_old_jobs()
            continue
        except Exception as e:
            print(f"❌ Assembly worker error: {e}")
            time.sleep(1)


def start_assembly_worker():
    """Start the background assembly worker"""
    worker_thread = threading.Thread(target=assembly_worker, daemon=True)
    worker_thread.start()
    print("🚀 Started background assembly worker")


def detect_ready_assemblies():
    """Detect chunks that are ready for assembly on startup"""
    try:
        chunks_dir = os.path.join(ROOT_DIR, ".chunks")
        if not os.path.exists(chunks_dir):
            return

        recovered_count = 0

        for file_id in os.listdir(chunks_dir):
            chunk_dir = os.path.join(chunks_dir, file_id)
            if not os.path.isdir(chunk_dir):
                continue

            try:
                # Skip if assembly is currently in progress
                protection_file = os.path.join(chunk_dir, ".assembling")
                if os.path.exists(protection_file):
                    print(f"🛡️ Skipping {file_id} - assembly protection active")
                    continue

                # Look for metadata file first
                metadata_file = os.path.join(chunk_dir, ".metadata")
                filename = f"recovered_file_{file_id}"
                dest_path = ""
                expected_chunks = None

                if os.path.exists(metadata_file):
                    try:
                        with open(metadata_file, "r") as f:
                            metadata = json.load(f)
                            filename = metadata.get("filename", filename)
                            dest_path = metadata.get("dest_path", dest_path)
                            expected_chunks = metadata.get("total_chunks")

                            print(
                                f"📋 Found metadata for {file_id}: {filename}, expected {expected_chunks} chunks"
                            )
                    except Exception as e:
                        print(f"⚠️ Error reading metadata for {file_id}: {e}")
                        continue

                # Use enhanced chunk verification
                try:
                    chunk_info = storage.verify_chunks_complete(
                        file_id, expected_chunks
                    )
                    total_chunks = chunk_info["total_chunks"]

                    print(
                        f"🔄 Found complete upload ready for assembly: {filename} ({total_chunks} chunks)"
                    )
                    assembly_queue.add_job(file_id, filename, dest_path, total_chunks)
                    recovered_count += 1

                except Exception as verify_error:
                    print(f"⚠️ Chunk verification failed for {file_id}: {verify_error}")
                    # Could cleanup incomplete uploads here if desired
                    continue

            except Exception as e:
                print(f"⚠️ Error checking chunks for {file_id}: {e}")
                continue

        if recovered_count > 0:
            print(
                f"🔄 Recovered {recovered_count} incomplete upload(s) for background assembly"
            )
        else:
            print("🔍 No incomplete uploads found ready for recovery")

    except Exception as e:
        print(f"⚠️ Error detecting ready assemblies: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# HLS Adaptive Streaming  (ffmpeg backend → Video.js frontend)
# ─────────────────────────────────────────────────────────────────────────────

_VIDEO_EXTS_HLS = frozenset(
    [
        "mp4",
        "webm",
        "mov",
        "m4v",
        "mkv",
        "avi",
        "wmv",
        "flv",
        "mpg",
        "mpeg",
        "m2ts",
        "mts",
        "3gp",
        "ogv",
        "ts",
    ]
)

# CRF-based quality ladder — no target bitrate, just a maxrate ceiling.
# Encoding uses -crf 18 (visually near-lossless) + -maxrate/-bufsize to cap
# runaway bitrates on complex scenes, same approach as YouTube.
#
# (name, target_height, maxrate, bufsize, audio_bitrate)
_HLS_BASE_PROFILES = [
    ("2160p", 2160, "40000k", "80000k", "192k"),  # 4K  — ~40 Mbps ceiling
    ("1440p", 1440, "24000k", "48000k", "192k"),  # 2K  — ~24 Mbps ceiling
    ("1080p", 1080, "12000k", "24000k", "192k"),  # matches YouTube Premium
    ("720p", 720, "7500k", "15000k", "128k"),
    ("480p", 480, "4000k", "8000k", "128k"),
    ("360p", 360, "1500k", "3000k", "96k"),
    ("240p", 240, "800k", "1600k", "64k"),
    ("144p", 144, "300k", "600k", "64k"),
]

# High-frame-rate variants — only added when source ≥ 48 fps, for 720p and above.
# Ceilings are ~50% higher than their SDR counterparts to accommodate extra frames.
# (name, target_height, maxrate, bufsize, audio_bitrate)
_HLS_HFR_PROFILES = [
    ("2160p60", 2160, "60000k", "120000k", "192k"),
    ("1440p60", 1440, "36000k", "72000k", "192k"),
    ("1080p60", 1080, "20000k", "40000k", "192k"),
    ("720p60", 720, "12000k", "24000k", "128k"),
]

_HLS_SEG_DURATION = 6  # seconds per HLS segment

# HLS cache root lives in the configured cache dir, not inside ROOT_DIR
from paths import get_hls_cache_dir as _get_hls_cache_dir


def _hls_cache_root() -> str:
    """Return (and create) the HLS cache directory from paths configuration."""
    return _get_hls_cache_dir(create=True)


def _hls_cache_key(full_path: str) -> str:
    """
    Stable 32-char hex key derived from (file_size, mtime).
    A renamed/moved file whose content is unchanged reuses the same key
    and avoids a full re-transcode.
    """
    try:
        st = os.stat(full_path)
        fingerprint = f"{st.st_size}:{st.st_mtime}"
    except OSError:
        fingerprint = f"{full_path}:0"
    return hashlib.md5(fingerprint.encode()).hexdigest()


def _hls_output_dir(cache_key: str) -> str:
    return os.path.join(_hls_cache_root(), cache_key)


_hls_status_lock = threading.Lock()


def _hls_read_status(cache_key: str) -> dict:
    f = os.path.join(_hls_output_dir(cache_key), ".status.json")
    if not os.path.exists(f):
        return {"status": "not_started"}
    try:
        with _hls_status_lock:
            with open(f, "r", encoding="utf-8") as fh:
                return json.load(fh)
    except Exception:
        return {"status": "unknown"}


def _hls_write_status(cache_key: str, data: dict):
    """
    Write .status.json under a threading lock.
    os.replace() is avoided because on Windows it raises PermissionError
    when another thread has the destination file open for reading at the
    same time (WinError 5).  A lock + direct overwrite is safe here because
    every reader also holds the same lock, so reads and writes never race.
    """
    d = _hls_output_dir(cache_key)
    os.makedirs(d, exist_ok=True)
    path = os.path.join(d, ".status.json")
    with _hls_status_lock:
        with open(path, "w", encoding="utf-8") as fh:
            json.dump(data, fh)


def _probe_video(file_path: str):
    """
    Return (width, height, has_audio, duration_secs, fps) via ffprobe.
    fps is the real/exact frame rate from r_frame_rate (e.g. 59.94, 60.0, 30.0).
    Falls back to (0, 0, False, 0.0, 0.0) on any error.
    """
    try:
        _ffmpeg_bin = _resolve_ffmpeg()
        _ffprobe_name = "ffprobe.exe" if os.name == "nt" else "ffprobe"
        ffprobe_bin = os.path.join(os.path.dirname(_ffmpeg_bin), _ffprobe_name)
        r = subprocess.run(
            [
                ffprobe_bin,
                "-v",
                "quiet",
                "-print_format",
                "json",
                "-show_streams",
                "-show_format",
                file_path,
            ],
            capture_output=True,
            text=True,
            timeout=30,
        )
        if r.returncode != 0:
            return 0, 0, False, 0.0, 0.0
        data = json.loads(r.stdout)
        streams = data.get("streams", [])
        fmt = data.get("format", {})
        w = h = 0
        has_audio = False
        fps = 0.0
        for s in streams:
            if s.get("codec_type") == "video":
                w = int(s.get("width") or 0)
                h = int(s.get("height") or 0)
                # r_frame_rate is exact rational e.g. "60000/1001" or "30/1"
                rfr = s.get("r_frame_rate") or s.get("avg_frame_rate") or "0/1"
                try:
                    num, den = rfr.split("/")
                    fps = float(int(num)) / float(int(den)) if int(den) != 0 else 0.0
                except Exception:
                    fps = 0.0
            elif s.get("codec_type") == "audio":
                has_audio = True
        try:
            duration = float(fmt.get("duration") or 0)
        except (TypeError, ValueError):
            duration = 0.0
        return w, h, has_audio, duration, fps
    except Exception:
        return 0, 0, False, 0.0, 0.0


def _resolve_ffmpeg() -> str:
    """Return the ffmpeg executable path, searching PATH then common Windows install dirs."""
    import shutil as _shutil

    found = _shutil.which("ffmpeg")
    if found:
        return found
    candidates = [
        os.path.expandvars(r"%LOCALAPPDATA%\Microsoft\WinGet\Packages"),
        os.path.expandvars(r"%ProgramFiles%\ffmpeg\bin"),
        os.path.expandvars(r"%ProgramFiles(x86)%\ffmpeg\bin"),
        r"C:\ProgramData\chocolatey\bin",
        os.path.expandvars(r"%USERPROFILE%\scoop\shims"),
        r"C:\ffmpeg\bin",
        r"C:\tools\ffmpeg\bin",
    ]
    for base in candidates:
        if not os.path.isdir(base):
            continue
        for root, dirs, files in os.walk(base):
            for fname in files:
                if fname.lower() in ("ffmpeg.exe", "ffmpeg"):
                    return os.path.join(root, fname)
            if root.count(os.sep) - base.count(os.sep) >= 5:
                dirs.clear()
    return "ffmpeg"  # last resort


def _ffmpeg_available() -> bool:
    """Quick check: is ffmpeg installed and executable?"""
    import shutil as _shutil

    bin_path = _resolve_ffmpeg()
    if bin_path == "ffmpeg" and not _shutil.which("ffmpeg"):
        return False
    try:
        subprocess.run(
            [bin_path, "-version"], capture_output=True, timeout=5, check=True
        )
        return True
    except Exception:
        return False


def _run_hls_transcode(file_path: str, cache_key: str):
    """
    Background daemon thread: transcode video → multi-quality HLS.

    Profile selection:
      • Standard profiles (144p→4K, capped at 30 fps) are included up to the
        source height.
      • HFR profiles (720p60→4K60) are added only when the source is ≥ 48 fps,
        for heights ≤ source height and ≥ 720.
      • Each standard profile applies an fps=30 cap via filter when the source
        is HFR; HFR profiles pass the source frame-rate through unchanged.
      • GOP size is set per-stream to 2 × effective_fps (2-second keyframe
        interval) so seeking stays accurate at any frame-rate.

    Writes live progress 0-100 to .status.json via -progress pipe:1 parsing.
    Survives frontend refresh — never interrupted by the browser.
    """
    output_dir = _hls_output_dir(cache_key)
    os.makedirs(output_dir, exist_ok=True)
    _hls_write_status(cache_key, {"status": "processing", "progress": 0})

    try:
        _, src_height, has_audio, duration_secs, src_fps = _probe_video(file_path)

        # Treat 48+ fps sources as HFR (covers both 50 Hz / PAL and 59.94/60 Hz)
        is_hfr = src_fps >= 48.0

        # ── Build active profile list ─────────────────────────────────────────
        # Each entry: (name, height, fps_cap, maxr, bufs, abr)
        #   fps_cap = 30    → standard profile; filter limits fps to 30 when HFR source
        #   fps_cap = None  → HFR profile; source frame-rate passes through unchanged
        profiles = []

        for name, h, maxr, bufs, abr in _HLS_BASE_PROFILES:
            if src_height == 0 or h <= src_height:
                profiles.append((name, h, 30, maxr, bufs, abr))

        if is_hfr:
            for name, h, maxr, bufs, abr in _HLS_HFR_PROFILES:
                if src_height == 0 or (h <= src_height and h >= 720):
                    profiles.append((name, h, None, maxr, bufs, abr))

        if not profiles:
            # Absolute fallback: lowest rung of the standard ladder
            name, h, maxr, bufs, abr = _HLS_BASE_PROFILES[-1]
            profiles = [(name, h, 30, maxr, bufs, abr)]

        n = len(profiles)

        for name, *_ in profiles:
            os.makedirs(os.path.join(output_dir, name), exist_ok=True)

        # ── Build filter_complex ──────────────────────────────────────────────
        # For standard profiles on an HFR source, append ",fps=fps=30" so the
        # 30-fps streams are correctly limited.  HFR profiles get no fps filter.
        splits = "".join(f"[vsp{i}]" for i in range(n))
        filter_parts = [f"[0:v]split={n}{splits}"]
        for i, (name, h, fps_cap, *_) in enumerate(profiles):
            fps_filter = (
                f",fps=fps={fps_cap}" if (fps_cap is not None and is_hfr) else ""
            )
            filter_parts.append(f"[vsp{i}]scale=-2:{h}{fps_filter}[vout{i}]")
        filter_complex = "; ".join(filter_parts)

        cmd = [
            _resolve_ffmpeg(),
            "-y",
            "-i",
            file_path,
            "-filter_complex",
            filter_complex,
        ]

        for i in range(n):
            cmd += ["-map", f"[vout{i}]"]
        if has_audio:
            for _ in range(n):
                cmd += ["-map", "0:a:0"]

        # ── Per-stream video encoder options (CRF mode) ──────────────────────
        # -crf 18 targets perceptual near-lossless quality; -maxrate/-bufsize
        # cap the bitrate ceiling so complex scenes don't explode, identical
        # to how YouTube's encoder pipeline works.
        for i, (name, h, fps_cap, maxr, bufs, abr) in enumerate(profiles):
            eff_fps = (fps_cap if fps_cap is not None else src_fps) or 30
            gop = max(48, int(round(eff_fps * 2)))
            cmd += [
                f"-c:v:{i}",
                "libx264",
                f"-crf:v:{i}",
                "22",
                f"-maxrate:v:{i}",
                maxr,
                f"-bufsize:v:{i}",
                bufs,
                f"-preset:v:{i}",
                "fast",
                f"-g:v:{i}",
                str(gop),
                f"-keyint_min:v:{i}",
                str(gop),
                f"-sc_threshold:v:{i}",
                "0",
            ]

        if has_audio:
            for i, (name, h, fps_cap, maxr, bufs, abr) in enumerate(profiles):
                cmd += [f"-c:a:{i}", "aac", f"-b:a:{i}", abr, "-ar", "48000"]

        if has_audio:
            vsm = " ".join(
                f"v:{i},a:{i},name:{name}" for i, (name, *_) in enumerate(profiles)
            )
        else:
            vsm = " ".join(
                f"v:{i},name:{name}" for i, (name, *_) in enumerate(profiles)
            )

        seg_tpl = os.path.join(output_dir, "%v", "seg%03d.ts")
        list_tpl = os.path.join(output_dir, "%v", "index.m3u8")

        cmd += [
            "-progress",
            "pipe:1",  # newline-delimited key=value → stdout
            "-nostats",  # suppress \r stats that block line iteration
            "-f",
            "hls",
            "-hls_time",
            str(_HLS_SEG_DURATION),
            "-hls_playlist_type",
            "vod",
            "-hls_flags",
            "independent_segments",
            "-var_stream_map",
            vsm,
            "-master_pl_name",
            "master.m3u8",
            "-hls_segment_filename",
            seg_tpl,
            list_tpl,
        ]

        profile_names = [p[0] for p in profiles]
        print(
            f"🎬 HLS transcode start: {os.path.basename(file_path)}  "
            f"src={src_height}p {'HFR({:.2f}fps)'.format(src_fps) if is_hfr else '{:.2f}fps'.format(src_fps)}  "
            f"profiles={profile_names}  duration={duration_secs:.1f}s",
            flush=True,
        )

        proc = subprocess.Popen(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            bufsize=1,
        )

        stderr_tail = []
        for line in proc.stdout:
            line = line.rstrip()
            if not line:
                continue
            if "=" in line:
                key, _, val = line.partition("=")
                key = key.strip()
                val = val.strip()
                print(f"[ffmpeg progress] {key}={val}", flush=True)
                # out_time_ms is in microseconds
                if key == "out_time_ms" and duration_secs > 0:
                    try:
                        pct = min(99, int(int(val) / 1_000_000 / duration_secs * 100))
                        _hls_write_status(
                            cache_key, {"status": "processing", "progress": pct}
                        )
                    except (ValueError, ZeroDivisionError):
                        pass
            else:
                print(f"[ffmpeg] {line}", flush=True)
                stderr_tail.append(line)
                if len(stderr_tail) > 50:
                    stderr_tail.pop(0)

        proc.wait()

        if proc.returncode == 0:
            _hls_write_status(
                cache_key,
                {
                    "status": "ready",
                    "progress": 100,
                    "profiles": [p[0] for p in profiles],
                },
            )
            print(f"\u2705 HLS transcode done: {cache_key[:8]}\u2026", flush=True)
        else:
            err_tail = "\n".join(stderr_tail)[-800:]
            _hls_write_status(cache_key, {"status": "error", "message": err_tail})
            print(
                f"\u274c HLS transcode failed (rc={proc.returncode}):\n{err_tail}",
                flush=True,
            )

    except subprocess.TimeoutExpired:
        _hls_write_status(
            cache_key, {"status": "error", "message": "Transcoding timed out"}
        )
        print(f"\u274c HLS transcode timed out: {cache_key[:8]}\u2026")
    except Exception as exc:
        _hls_write_status(cache_key, {"status": "error", "message": str(exc)})
        print(f"\u274c HLS transcode error: {exc}")


@app.route("/hls_start/<path:video_path>")
@login_required
def hls_start(video_path):
    """
    Kick off HLS transcoding (idempotent).  Returns:
      hls_available: false  — ffmpeg not installed; frontend shows Play Raw only
      status: processing    — transcode running; frontend shows progress bar
      status: ready         — segments ready;    frontend shows Stream HLS button
    """
    if not storage.is_safe_path(video_path):
        return jsonify({"error": "Invalid path"}), 400

    full_path = os.path.join(ROOT_DIR, video_path)
    if not os.path.exists(full_path) or os.path.isdir(full_path):
        return jsonify({"error": "File not found"}), 404

    ext = video_path.rsplit(".", 1)[-1].lower() if "." in video_path else ""
    if ext not in _VIDEO_EXTS_HLS:
        return jsonify({"hls_available": False, "reason": "unsupported_format"})

    # Skip HLS for small web-native files — play raw is fine and saves CPU/disk.
    # Non-web-native formats (mkv, avi, wmv, etc.) always get HLS regardless of
    # size because the browser cannot decode them natively.
    # Both thresholds are configurable via config.py / server_config.json.
    _web_native = {"mp4", "webm", "mov", "m4v", "ts"}
    if HLS_MIN_SIZE > 0 and ext in _web_native and ext not in HLS_FORCE_FORMATS:
        try:
            file_size = os.path.getsize(full_path)
        except OSError:
            file_size = 0
        if file_size < HLS_MIN_SIZE:
            return jsonify(
                {
                    "hls_available": False,
                    "reason": "file_too_small",
                    "file_size": file_size,
                    "min_size": HLS_MIN_SIZE,
                }
            )

    if not _ffmpeg_available():
        print(
            "\u26a0\ufe0f  ffmpeg not found — HLS unavailable, client will use raw playback"
        )
        return jsonify({"hls_available": False, "reason": "ffmpeg_not_installed"})

    cache_key = _hls_cache_key(full_path)
    status = _hls_read_status(cache_key)

    # Guard stale "ready" when hls cache dir was wiped externally
    if status.get("status") == "ready":
        master = os.path.join(_hls_output_dir(cache_key), "master.m3u8")
        if os.path.exists(master):
            return jsonify({"hls_available": True, "cache_key": cache_key, **status})
        print(
            f"\u26a0\ufe0f  HLS stale cache — re-transcoding: {cache_key[:8]}\u2026",
            flush=True,
        )
        # fall through

    if status.get("status") == "processing":
        return jsonify({"hls_available": True, "cache_key": cache_key, **status})

    # Spawn background daemon thread — survives frontend refresh
    threading.Thread(
        target=_run_hls_transcode, args=(full_path, cache_key), daemon=True
    ).start()

    _hls_write_status(cache_key, {"status": "processing", "progress": 0})
    return jsonify(
        {
            "hls_available": True,
            "cache_key": cache_key,
            "status": "processing",
            "progress": 0,
        }
    )


@app.route("/hls_status/<cache_key>")
@login_required
def hls_status_route(cache_key):
    """Poll transcoding status.  Returns {status, progress 0-100, profiles?}."""
    if not re.fullmatch(r"[a-f0-9]{32}", cache_key):
        return jsonify({"error": "Invalid key"}), 400
    return jsonify(_hls_read_status(cache_key))


@app.route("/hls_files/<cache_key>/<path:hls_path>")
@login_required
def hls_files(cache_key, hls_path):
    """Serve HLS master/sub-playlists (.m3u8) and TS segments (.ts)."""
    if not re.fullmatch(r"[a-f0-9]{32}", cache_key):
        return "Invalid key", 400
    if ".." in hls_path or hls_path.startswith("/"):
        return "Forbidden", 403
    _, ext = os.path.splitext(hls_path)
    if ext.lower() not in (".m3u8", ".ts"):
        return "Forbidden", 403
    output_dir = _hls_output_dir(cache_key)
    file_path = os.path.normpath(os.path.join(output_dir, hls_path))
    if not file_path.startswith(os.path.abspath(output_dir)):
        return "Path traversal", 403
    if not os.path.exists(file_path):
        return "Not found", 404
    if ext.lower() == ".m3u8":
        return send_file(
            file_path,
            mimetype="application/vnd.apple.mpegurl",
            max_age=0,
            conditional=False,
        )
    return send_file(file_path, mimetype="video/mp2t", max_age=3600)


# ─────────────────────────────────────────────────────────────────────────────
# End HLS Adaptive Streaming
# ─────────────────────────────────────────────────────────────────────────────


# Initialize cleanup on startup
def initialize_cleanup():
    """Initialize all cleanup processes"""
    print("🧹 Initializing cleanup systems...")

    # Start enhanced cleanup schedulers
    start_enhanced_cleanup_scheduler()
    start_orphan_cleanup_scheduler()

    # Start assembly worker
    start_assembly_worker()

    # Do an initial aggressive cleanup on startup
    try:
        print("🧹 Running startup cleanup...")

        # Get active assembly jobs to protect them (should be none on startup)
        active_assembly_jobs = get_protected_files()

        storage.cleanup_old_chunks(
            max_age_hours=0.1, protected_files=active_assembly_jobs
        )  # Clean chunks older than 6 minutes
        chunk_tracker.cleanup_orphaned_chunks()
        print("✅ Startup cleanup completed")
    except Exception as e:
        print(f"⚠️ Warning: Startup cleanup failed: {e}")

    # Check for any existing chunks that are ready for assembly
    print("🔍 Checking for incomplete uploads ready for assembly...")
    detect_ready_assemblies()


# Initialize cleanup when app starts
initialize_cleanup()


if __name__ == "__main__":
    print(f"🚀 Starting Enhanced Cloudinator FTP Server on port {PORT}")
    print(f"📁 Root directory: {os.path.abspath(ROOT_DIR)}")
    print(f"🔧 Chunked uploads: {'Enabled' if ENABLE_CHUNKED_UPLOADS else 'Disabled'}")
    print(f"📦 Chunk size: {CHUNK_SIZE // (1024*1024)}MB")
    print("✨ Enhanced Features:")
    print("   • Smart progress tracking with speed/ETA")
    print("   • Multi-file selection with bulk operations")
    print("   • Advanced chunk cleanup system")
    print("   • Session-based upload tracking")
    print("   • Orphaned chunk detection and cleanup")
    print("   • Real-time cleanup on page refresh")
    print("   • Background file assembly with status tracking")
    print("🧹 Cleanup Schedule:")
    print("   • Every 5 minutes: Orphaned chunks cleanup")
    print("   • Every 15 minutes: Stale chunks cleanup (1+ hours)")
    print("   • Every 1 hour: Full cleanup (24+ hours)")
    print("   • On page load: Request-based cleanup")
    print("   • On logout: Session cleanup")
    print("🔄 Assembly System:")
    print("   • Background worker processes file assembly")
    print("   • Real-time status updates via API")
    print("   • Resume capability after page refresh")
    print("   • Automatic recovery of incomplete uploads")

    app.run(host="0.0.0.0", port=PORT, debug=False)
